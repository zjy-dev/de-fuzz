"""Theme / layout primitives for the DeFuzz 开题答辩 deck.

Centralizes palette, slide geometry, and low-level python-pptx helpers so the
per-slide builders stay declarative.
"""
from __future__ import annotations

from pathlib import Path
from typing import Iterable, Sequence

from lxml import etree
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

# ─────────────────────────────────────────────────────────────
# Palette (SEU 深绿 + 黄强调 + 橘 motivation 框)
# ─────────────────────────────────────────────────────────────
GREEN_DARK = RGBColor(0x1B, 0x5E, 0x20)      # primary
GREEN_MID = RGBColor(0x2E, 0x7D, 0x32)
GREEN_LIGHT = RGBColor(0xE8, 0xF5, 0xE9)
GREEN_PALE = RGBColor(0xC8, 0xE6, 0xC9)
YELLOW_ACCENT = RGBColor(0xFF, 0xC1, 0x07)
YELLOW_LIGHT = RGBColor(0xFF, 0xF8, 0xE1)
ORANGE_DARK = RGBColor(0xE6, 0x51, 0x00)
ORANGE_MID = RGBColor(0xFF, 0x6F, 0x00)
ORANGE_LIGHT = RGBColor(0xFF, 0xF3, 0xE0)
RED_DARK = RGBColor(0xC6, 0x28, 0x28)
RED_LIGHT = RGBColor(0xFF, 0xEB, 0xEE)
BLUE_DARK = RGBColor(0x15, 0x65, 0xC0)
BLUE_LIGHT = RGBColor(0xE3, 0xF2, 0xFD)
GRAY_50 = RGBColor(0xFA, 0xFA, 0xFA)
GRAY_100 = RGBColor(0xF5, 0xF5, 0xF5)
GRAY_300 = RGBColor(0xE0, 0xE0, 0xE0)
GRAY_500 = RGBColor(0x9E, 0x9E, 0x9E)
GRAY_600 = RGBColor(0x66, 0x66, 0x66)
GRAY_700 = RGBColor(0x55, 0x55, 0x55)
GRAY_800 = RGBColor(0x33, 0x33, 0x33)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)

# ─────────────────────────────────────────────────────────────
# Geometry (16:9, 10" × 5.625")
# ─────────────────────────────────────────────────────────────
SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)

HEADER_H = Inches(0.55)
FOOTER_H = Inches(0.08)
CONTENT_PAD_X = Inches(0.35)
CONTENT_PAD_Y = Inches(0.18)
CONTENT_TOP = HEADER_H + CONTENT_PAD_Y
CONTENT_W = SLIDE_W - 2 * CONTENT_PAD_X
CONTENT_H = SLIDE_H - HEADER_H - FOOTER_H - 2 * CONTENT_PAD_Y

# Fonts — 全 deck 统一使用微软雅黑 (中英文同字体, 便于跨平台一致渲染)
LATIN_FONT = "Microsoft YaHei"
EA_FONT = "Microsoft YaHei"

# Workspace paths
WORKSPACE = Path(__file__).resolve().parent
LOGO_PATH = WORKSPACE / "seu-logo.png"
OUTPUT_PATH = WORKSPACE / "开题答辩-DeFuzz.pptx"


# ─────────────────────────────────────────────────────────────
# Page counter (扉页不占页码, 内容页 next() 递增)
# ─────────────────────────────────────────────────────────────
class PageCounter:
    """页码计数器: 扉页 builder 不调用, 内容 builder 调 next() 取页号."""

    def __init__(self) -> None:
        self.n = 0

    def next(self) -> int:
        self.n += 1
        return self.n


# ─────────────────────────────────────────────────────────────
# Low-level XML helper for font + east-asian typeface
# ─────────────────────────────────────────────────────────────
def _apply_font_to_run(run, *, name: str = LATIN_FONT, ea_name: str = EA_FONT) -> None:
    """Force Latin + East Asian typeface so CJK characters render reliably."""
    rPr = run._r.get_or_add_rPr()
    for tag, typeface in ((qn("a:latin"), name), (qn("a:ea"), ea_name)):
        el = rPr.find(tag)
        if el is None:
            el = etree.SubElement(rPr, tag)
        el.set("typeface", typeface)


# ─────────────────────────────────────────────────────────────
# Shape primitives
# ─────────────────────────────────────────────────────────────
def add_rect(
    slide,
    x,
    y,
    w,
    h,
    *,
    fill: RGBColor | None,
    line: RGBColor | None = None,
    line_width: float = 0.75,
    shape=MSO_SHAPE.RECTANGLE,
):
    """Rectangle (or rounded) with solid fill. `fill=None` ⇒ no fill."""
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_width)
    sp.shadow.inherit = False
    return sp


def add_line(slide, x1, y1, x2, y2, *, color: RGBColor, width: float = 1.0):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)  # 1 = straight line
    line.line.color.rgb = color
    line.line.width = Pt(width)
    return line


# ─────────────────────────────────────────────────────────────
# Text primitives
# ─────────────────────────────────────────────────────────────
def _set_paragraph_align(p, align: PP_ALIGN | None):
    if align is not None:
        p.alignment = align


def add_text(
    slide,
    x,
    y,
    w,
    h,
    lines: Sequence[str | dict],
    *,
    font_size: float = 12,
    bold: bool = False,
    italic: bool = False,
    color: RGBColor = GRAY_800,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
    line_spacing: float | None = None,
    margin: float = 0.0,
    font_name: str = LATIN_FONT,
    ea_name: str = EA_FONT,
):
    """Add a text box. `lines` is a list; each element is either:
    - a plain string → a paragraph with default style, or
    - a dict `{text, size?, bold?, italic?, color?, align?, space_after?, runs?}`
      where `runs` is a list of `{text, size?, bold?, italic?, color?}` for mixed
      in-line styles. A paragraph with `runs` ignores `text`.
    """
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    m = Emu(int(margin))
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = m

    for i, item in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing is not None:
            p.line_spacing = line_spacing

        if isinstance(item, str):
            runs = [{"text": item}]
            p_opts: dict = {}
        else:
            runs = item.get("runs") or [{
                "text": item.get("text", ""),
                "size": item.get("size"),
                "bold": item.get("bold"),
                "italic": item.get("italic"),
                "color": item.get("color"),
            }]
            p_opts = item
            _set_paragraph_align(p, item.get("align"))
            if item.get("space_after") is not None:
                p.space_after = Pt(item["space_after"])
            if item.get("space_before") is not None:
                p.space_before = Pt(item["space_before"])
            if item.get("line_spacing") is not None:
                p.line_spacing = item["line_spacing"]

        # clear the auto-generated first run before appending
        if i == 0:
            for r in list(p.runs):
                r._r.getparent().remove(r._r)

        for run_spec in runs:
            r = p.add_run()
            r.text = run_spec.get("text", "")
            rf = r.font
            rf.size = Pt(run_spec.get("size") or p_opts.get("size") or font_size)
            rf.bold = bool(run_spec.get("bold") if run_spec.get("bold") is not None else p_opts.get("bold", bold))
            rf.italic = bool(run_spec.get("italic") if run_spec.get("italic") is not None else p_opts.get("italic", italic))
            rf.color.rgb = run_spec.get("color") or p_opts.get("color") or color
            rf.name = font_name
            _apply_font_to_run(r, name=font_name, ea_name=ea_name)
    return tb


# ─────────────────────────────────────────────────────────────
# Composite: header + footer for content slides
# ─────────────────────────────────────────────────────────────
def add_page_frame(slide, title: str, page_no: int | None = None):
    """Draw green header bar with title + yellow footer bar + page number."""
    # Header
    add_rect(slide, 0, 0, SLIDE_W, HEADER_H, fill=GREEN_DARK)
    add_text(
        slide,
        Inches(0.35),
        Inches(0.05),
        SLIDE_W - Inches(0.7),
        HEADER_H - Inches(0.05),
        [title],
        font_size=20,
        bold=True,
        color=WHITE,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    # Footer accent bar
    add_rect(slide, 0, SLIDE_H - FOOTER_H, SLIDE_W, FOOTER_H, fill=YELLOW_ACCENT)
    # Page number
    if page_no is not None:
        add_text(
            slide,
            SLIDE_W - Inches(0.8),
            SLIDE_H - FOOTER_H - Inches(0.28),
            Inches(0.7),
            Inches(0.22),
            [str(page_no)],
            font_size=9,
            color=GRAY_600,
            align=PP_ALIGN.RIGHT,
        )


# ─────────────────────────────────────────────────────────────
# Small section header ribbon inside a slide (yellow bar + green title)
# ─────────────────────────────────────────────────────────────
def add_section_title(slide, x, y, w, text: str, *, font_size: float = 13):
    bar_w = Inches(0.06)
    add_rect(slide, x, y + Emu(int(Pt(2))), bar_w, Pt(font_size + 2), fill=YELLOW_ACCENT)
    add_text(
        slide,
        x + bar_w + Inches(0.08),
        y,
        w - bar_w - Inches(0.08),
        Pt(font_size + 6),
        [text],
        font_size=font_size,
        bold=True,
        color=GREEN_DARK,
        anchor=MSO_ANCHOR.MIDDLE,
    )


# ─────────────────────────────────────────────────────────────
# Utility: arrow between path boxes
# ─────────────────────────────────────────────────────────────
def add_arrow(slide, x, y, w, h, *, fill: RGBColor, line: RGBColor | None = None):
    """Right-pointing arrow shape."""
    sp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
    sp.shadow.inherit = False
    return sp


__all__ = [
    "GREEN_DARK", "GREEN_MID", "GREEN_LIGHT", "GREEN_PALE",
    "YELLOW_ACCENT", "YELLOW_LIGHT",
    "ORANGE_DARK", "ORANGE_MID", "ORANGE_LIGHT",
    "RED_DARK", "RED_LIGHT",
    "BLUE_DARK", "BLUE_LIGHT",
    "GRAY_50", "GRAY_100", "GRAY_300", "GRAY_500", "GRAY_600", "GRAY_700", "GRAY_800",
    "WHITE", "BLACK",
    "SLIDE_W", "SLIDE_H", "HEADER_H", "FOOTER_H",
    "CONTENT_PAD_X", "CONTENT_PAD_Y", "CONTENT_TOP", "CONTENT_W", "CONTENT_H",
    "LOGO_PATH", "OUTPUT_PATH",
    "MSO_SHAPE", "PP_ALIGN", "MSO_ANCHOR",
    "Inches", "Pt", "Emu",
    "PageCounter",
    "add_rect", "add_line", "add_text", "add_page_frame", "add_section_title",
    "add_arrow",
]
