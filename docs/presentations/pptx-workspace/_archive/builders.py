"""Per-slide builders for the DeFuzz 开题答辩 deck.

Each function takes a `pres` (Presentation) and returns the added slide.
Coordinates use Inches/Pt via the helpers in `theme`.
"""
from __future__ import annotations

from pptx.util import Inches, Pt

from theme import (
    BLACK,
    BLUE_DARK,
    BLUE_LIGHT,
    CONTENT_H,
    CONTENT_PAD_X,
    CONTENT_PAD_Y,
    CONTENT_TOP,
    CONTENT_W,
    FOOTER_H,
    GRAY_100,
    GRAY_300,
    GRAY_500,
    GRAY_600,
    GRAY_700,
    GRAY_800,
    GREEN_DARK,
    GREEN_LIGHT,
    GREEN_MID,
    GREEN_PALE,
    HEADER_H,
    LOGO_PATH,
    MSO_ANCHOR,
    MSO_SHAPE,
    ORANGE_DARK,
    ORANGE_LIGHT,
    ORANGE_MID,
    PP_ALIGN,
    RED_DARK,
    RED_LIGHT,
    SLIDE_H,
    SLIDE_W,
    WHITE,
    YELLOW_ACCENT,
    YELLOW_LIGHT,
    add_arrow,
    add_line,
    add_page_frame,
    add_rect,
    add_section_title,
    add_text,
)

BLANK_LAYOUT = 6  # python-pptx default blank layout index


# ─────────────────────────────────────────────────────────────
# Slide 1 — Title
# ─────────────────────────────────────────────────────────────
def slide_01_title(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[BLANK_LAYOUT])
    # Green + yellow accent stripes
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.12), fill=GREEN_DARK)
    add_rect(slide, 0, SLIDE_H - Inches(0.08), SLIDE_W, Inches(0.08), fill=YELLOW_ACCENT)

    # Logo centered
    logo_w = Inches(0.95)
    slide.shapes.add_picture(
        str(LOGO_PATH),
        (SLIDE_W - logo_w) / 2,
        Inches(1.1),
        width=logo_w,
        height=logo_w,
    )

    # Main title
    add_text(
        slide,
        Inches(0.5), Inches(2.25), Inches(9), Inches(0.9),
        ["DeFuzz: LLM 驱动的编译器软件防御机制模糊测试"],
        font_size=30, bold=True, color=GREEN_DARK,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )
    # Subtitle
    add_text(
        slide,
        Inches(0.5), Inches(3.3), Inches(9), Inches(0.5),
        ["硕士学位论文开题答辩"],
        font_size=18, color=GRAY_700,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )
    # Info line
    add_text(
        slide,
        Inches(0.5), Inches(4.35), Inches(9), Inches(0.4),
        ["答辩人: XXX    导师: XXX    2026 年 2 月"],
        font_size=13, color=GRAY_600,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )
    return slide


# ─────────────────────────────────────────────────────────────
# Slide 2 — 研究背景
# ─────────────────────────────────────────────────────────────
def slide_02_background(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[BLANK_LAYOUT])
    add_page_frame(slide, "1. 研究背景 — 编译器软件防御机制", page_no=1)

    y = CONTENT_TOP
    add_section_title(
        slide, CONTENT_PAD_X, y, CONTENT_W,
        "编译器为用户代码生成的软件防御机制, 是系统安全的关键防线",
    )

    y_cards = y + Inches(0.45)
    col_w = (CONTENT_W - Inches(0.25)) / 2
    left_x = CONTENT_PAD_X
    right_x = left_x + col_w + Inches(0.25)
    cards_h = CONTENT_H - Inches(0.5)

    # Left: 主流防御机制
    add_rect(slide, left_x, y_cards, col_w, Inches(0.32), fill=GREEN_DARK)
    add_text(
        slide,
        left_x + Inches(0.1), y_cards, col_w - Inches(0.2), Inches(0.32),
        ["主流编译器防御机制 (DeFuzz 覆盖面)"],
        font_size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE,
    )

    mech_items = [
        ("Stack Canary", "栈保护 canary, 检测缓冲区溢出覆盖返回地址"),
        ("FORTIFY_SOURCE", "对 memcpy / strcpy / strcat 等注入 _chk 变体做边界检查"),
        ("CET-IBT / BTI / PAC", "控制流完整性: 间接跳转 landing pad + 指针鉴权"),
        ("SafeStack / ShadowStack / StackClash", "栈硬化: 分离敏感栈 + 探针防跨越"),
    ]
    item_h = Inches(0.62)
    item_y = y_cards + Inches(0.38)
    for name, desc in mech_items:
        add_rect(
            slide, left_x, item_y, col_w, item_h,
            fill=GREEN_LIGHT, line=GREEN_PALE,
        )
        # left yellow accent
        add_rect(slide, left_x, item_y, Inches(0.08), item_h, fill=YELLOW_ACCENT)
        add_text(
            slide,
            left_x + Inches(0.18), item_y + Inches(0.06),
            col_w - Inches(0.25), Inches(0.28),
            [name], font_size=12, bold=True, color=GREEN_DARK,
        )
        add_text(
            slide,
            left_x + Inches(0.18), item_y + Inches(0.32),
            col_w - Inches(0.25), item_h - Inches(0.32),
            [desc], font_size=10, color=GRAY_700, line_spacing=1.1,
        )
        item_y += item_h + Inches(0.08)

    # Right: 问题
    add_rect(slide, right_x, y_cards, col_w, Inches(0.32), fill=RED_DARK)
    add_text(
        slide,
        right_x + Inches(0.1), y_cards, col_w - Inches(0.2), Inches(0.32),
        ["问题: 静默的逻辑漏洞"],
        font_size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE,
    )

    problem_items = [
        (
            "无崩溃 · 无报错",
            "防御机制失效时程序表面行为正常, 漏洞潜伏在 compiler backend 与 libc runtime 之间",
        ),
        (
            "传统 crash-oracle 失效",
            "AFL / libFuzzer 依赖段错误、assertion、ASan 报错, 对 ' 编译器没按规范生成防护代码 ' 无感",
        ),
        (
            "安全影响巨大",
            "单个 bug 可使整类二进制的安全保证沦为纸糊, 典型代价参见 CVE-2023-4039 (AArch64 canary 绕过)",
        ),
    ]
    p_item_h = Inches(0.85)
    p_item_y = y_cards + Inches(0.38)
    for title, desc in problem_items:
        add_rect(
            slide, right_x, p_item_y, col_w, p_item_h,
            fill=RED_LIGHT, line=RED_DARK, line_width=0.5,
        )
        add_text(
            slide,
            right_x + Inches(0.15), p_item_y + Inches(0.08),
            col_w - Inches(0.3), Inches(0.3),
            [title], font_size=12, bold=True, color=RED_DARK,
        )
        add_text(
            slide,
            right_x + Inches(0.15), p_item_y + Inches(0.38),
            col_w - Inches(0.3), p_item_h - Inches(0.4),
            [desc], font_size=10, color=GRAY_800, line_spacing=1.2,
        )
        p_item_y += p_item_h + Inches(0.08)

    return slide


# ─────────────────────────────────────────────────────────────
# Slide 3 — 研究现状
# ─────────────────────────────────────────────────────────────
def slide_03_status(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[BLANK_LAYOUT])
    add_page_frame(slide, "2. 研究现状 — LLM 驱动的编译器 Fuzzing", page_no=2)

    y = CONTENT_TOP
    add_section_title(
        slide, CONTENT_PAD_X, y, CONTENT_W,
        "近年 LLM 引入编译器 fuzzing, 但聚焦优化正确性, 未覆盖软件防御机制",
    )

    y_cards = y + Inches(0.45)
    col_w = (CONTENT_W - Inches(0.25)) / 2
    left_x = CONTENT_PAD_X
    right_x = left_x + col_w + Inches(0.25)
    card_h = Inches(2.55)

    works = [
        {
            "x": left_x,
            "title": "WhiteFox (OSDI'24)",
            "tag": "LLM 引导路径触发",
            "bullets": [
                "用 LLM 阅读优化 pass 源码, 生成触发特定优化的 C 程序",
                "擅长 miscompile 等 observable bug",
                "依赖手工编写优化规则, 覆盖面有限",
            ],
        },
        {
            "x": right_x,
            "title": "HLPFuzz (ICSE'25)",
            "tag": "LLM 变异 + 覆盖率引导",
            "bullets": [
                "把 LLM 作为变异算子与覆盖率 fuzzer 混合",
                "聚焦语义 bug (crash / miscompile), 以 crash oracle 为主",
                "缺少 invariant 级别的防御机制语义契约",
            ],
        },
    ]
    for w in works:
        x = w["x"]
        add_rect(slide, x, y_cards, col_w, Inches(0.42), fill=GREEN_DARK)
        add_text(
            slide,
            x + Inches(0.15), y_cards, col_w - Inches(0.2), Inches(0.42),
            [w["title"]],
            font_size=14, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE,
        )
        add_rect(
            slide, x, y_cards + Inches(0.42), col_w, card_h - Inches(0.42),
            fill=GRAY_100, line=GRAY_300,
        )
        add_text(
            slide,
            x + Inches(0.2), y_cards + Inches(0.52),
            col_w - Inches(0.4), Inches(0.28),
            [w["tag"]],
            font_size=11, bold=True, italic=True, color=GREEN_DARK,
        )
        bullets = [{"text": "• " + b, "space_after": 4} for b in w["bullets"]]
        add_text(
            slide,
            x + Inches(0.2), y_cards + Inches(0.88),
            col_w - Inches(0.4), card_h - Inches(1.0),
            bullets,
            font_size=11, color=GRAY_800, line_spacing=1.25,
        )

    # Limitation banner at bottom
    lim_y = y_cards + card_h + Inches(0.15)
    lim_h = Inches(0.72)
    add_rect(
        slide, CONTENT_PAD_X, lim_y, CONTENT_W, lim_h,
        fill=ORANGE_LIGHT, line=ORANGE_DARK, line_width=1.0,
    )
    add_rect(slide, CONTENT_PAD_X, lim_y, Inches(0.08), lim_h, fill=ORANGE_DARK)
    add_text(
        slide,
        CONTENT_PAD_X + Inches(0.2), lim_y + Inches(0.08),
        CONTENT_W - Inches(0.3), Inches(0.3),
        ["共同局限"],
        font_size=12, bold=True, color=ORANGE_DARK,
    )
    add_text(
        slide,
        CONTENT_PAD_X + Inches(0.2), lim_y + Inches(0.36),
        CONTENT_W - Inches(0.3), lim_h - Inches(0.38),
        [
            "现有工作均无针对 '软件防御机制是否被编译器正确生成并维持有效' 的形式化 oracle —— "
            "这恰是编译器安全链路上最脆弱、也最难自动化检测的环节。"
        ],
        font_size=11, color=GRAY_800, line_spacing=1.25,
    )

    return slide


# ─────────────────────────────────────────────────────────────
# Slide 4 — Motivation (Fortify-source bug)
# ─────────────────────────────────────────────────────────────
def slide_04_motivation(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[BLANK_LAYOUT])
    add_page_frame(slide, "3. 研究动机 — 编译器自身的防御实现仍在出错", page_no=3)

    y = CONTENT_TOP
    add_section_title(
        slide, CONTENT_PAD_X, y, CONTENT_W,
        "真实案例: FORTIFY_SOURCE 在 GCC 的优化路径上静默失效",
    )

    intro_y = y + Inches(0.42)
    add_text(
        slide,
        CONTENT_PAD_X, intro_y, CONTENT_W, Inches(0.45),
        [
            "FORTIFY_SOURCE 把 strcpy / strcat 等函数替换成 __*_chk 变体, 运行时比对目标缓冲区容量与写入长度 —— "
            "一旦 GCC 改写错了 _chk 参数, 这层保护就默默失效。"
        ],
        font_size=11, color=GRAY_800, line_spacing=1.3,
    )

    # Orange path box
    path_y = intro_y + Inches(0.55)
    path_h = Inches(1.5)
    add_rect(
        slide, CONTENT_PAD_X, path_y, CONTENT_W, path_h,
        fill=ORANGE_LIGHT, line=ORANGE_DARK, line_width=1.5,
    )
    add_text(
        slide,
        CONTENT_PAD_X + Inches(0.2), path_y + Inches(0.08),
        CONTENT_W - Inches(0.4), Inches(0.3),
        ["GCC 对 strcat 调用的改写链 (FORTIFY_SOURCE=2)"],
        font_size=11, bold=True, color=ORANGE_DARK,
    )

    # Three boxes + two arrows inside the orange frame (edge-to-edge, centered)
    box_w = Inches(2.55)
    box_h = Inches(0.7)
    arrow_w = Inches(0.35)
    arrow_h = Inches(0.3)
    group_w = box_w * 3 + arrow_w * 2
    inner_left = CONTENT_PAD_X + Inches(0.2)
    inner_w = CONTENT_W - Inches(0.4)
    assert inner_w > group_w, f"group too wide: {group_w} > {inner_w}"
    bx = inner_left + (inner_w - group_w) / 2
    by = path_y + Inches(0.48)

    steps = [
        ("strcat(dst, src)", "用户原始调用", WHITE, GREEN_DARK, GREEN_DARK),
        (
            "__strcat_chk(dst, src, objsize)",
            "安全改写 (FORTIFY wrapper)",
            GREEN_LIGHT,
            GREEN_DARK,
            GREEN_DARK,
        ),
        (
            "__stpcpy_chk(dst+strlen(dst), src, objsize)",
            "性能改写 (tree-ssa-strlen)",
            RED_LIGHT,
            RED_DARK,
            RED_DARK,
        ),
    ]
    for idx, (label, sub, fill, border, text_col) in enumerate(steps):
        add_rect(
            slide, bx, by, box_w, box_h,
            fill=fill, line=border, line_width=1.0,
        )
        add_text(
            slide,
            bx + Inches(0.05), by + Inches(0.05),
            box_w - Inches(0.1), Inches(0.36),
            [label],
            font_size=9.5, bold=True, color=text_col, align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        add_text(
            slide,
            bx + Inches(0.05), by + Inches(0.4),
            box_w - Inches(0.1), Inches(0.28),
            [sub],
            font_size=8.5, italic=True, color=GRAY_700, align=PP_ALIGN.CENTER,
        )
        if idx < 2:
            add_arrow(
                slide,
                bx + box_w,
                by + (box_h - arrow_h) / 2,
                arrow_w, arrow_h,
                fill=ORANGE_DARK,
            )
        bx += box_w + arrow_w

    # Caption below the path
    cap_y = path_y + path_h - Inches(0.32)
    add_text(
        slide,
        CONTENT_PAD_X + Inches(0.2), cap_y,
        CONTENT_W - Inches(0.4), Inches(0.3),
        [
            {
                "runs": [
                    {"text": "第三个参数 objsize ", "bold": True, "color": RED_DARK, "size": 10},
                    {"text": "应为 ", "size": 10},
                    {"text": "objsize - strlen(dst)", "bold": True, "color": RED_DARK, "size": 10},
                    {"text": " —— 优化路径没有调整, bound 仍按原缓冲区大小, 保护静默退化。", "size": 10},
                ],
                "align": PP_ALIGN.CENTER,
            }
        ],
        font_size=10, color=GRAY_800,
    )

    # Bottom status strip
    st_y = path_y + path_h + Inches(0.18)
    st_h = Inches(0.58)
    add_rect(
        slide, CONTENT_PAD_X, st_y, CONTENT_W, st_h,
        fill=GREEN_LIGHT, line=GREEN_DARK, line_width=0.75,
    )
    add_rect(slide, CONTENT_PAD_X, st_y, Inches(0.08), st_h, fill=GREEN_DARK)
    add_text(
        slide,
        CONTENT_PAD_X + Inches(0.2), st_y + Inches(0.04),
        CONTENT_W - Inches(0.3), Inches(0.26),
        ["上游状态"],
        font_size=11, bold=True, color=GREEN_DARK,
    )
    add_text(
        slide,
        CONTENT_PAD_X + Inches(0.2), st_y + Inches(0.3),
        CONTENT_W - Inches(0.3), st_h - Inches(0.32),
        [
            "已被 GCC 确认并 patch 进 master 分支, gcc 16.1+ 已修复; "
            "缺陷自 gcc 4.7 起存在约 14 年, 社区计划 backport 至 gcc 13 及以上。"
        ],
        font_size=10, color=GRAY_800, line_spacing=1.25,
    )
    return slide


# ─────────────────────────────────────────────────────────────
# Slide 5 — 研究目标与内容 (合并页)
# ─────────────────────────────────────────────────────────────
def slide_05_goals(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[BLANK_LAYOUT])
    add_page_frame(slide, "4. 研究目标与内容", page_no=4)

    y = CONTENT_TOP

    # Top: 3 goal cards
    add_section_title(slide, CONTENT_PAD_X, y, CONTENT_W, "研究目标 (Goals)")
    goals_y = y + Inches(0.4)
    goal_h = Inches(1.5)
    goal_gap = Inches(0.15)
    goal_w = (CONTENT_W - goal_gap * 2) / 3
    goals = [
        (
            "G1 · 抽取不变量",
            "系统化提取各防御机制的 invariants",
            ["来源: 源码注释 · 官方文档 · 已确认历史 bug", "24 份 invariant survey 作为形式化依据"],
        ),
        (
            "G2 · 构建预言机",
            "把每条 invariant 落成可执行 Checker",
            ["聚合为 MechanismOracle (多阶段调度)", "跨 ISA 统一接口, 支持正 / 负控"],
        ),
        (
            "G3 · 覆盖率驱动 Fuzzing",
            "以防御机制源码覆盖率为核心信号",
            ["CFG-guided target 选择 + LLM 种子变异", "发现静默逻辑漏洞 (GCC / LLVM)"],
        ),
    ]
    gx = CONTENT_PAD_X
    for title, subtitle, bullets in goals:
        add_rect(
            slide, gx, goals_y, goal_w, goal_h,
            fill=GREEN_LIGHT, line=GREEN_DARK, line_width=0.75,
        )
        add_rect(slide, gx, goals_y, goal_w, Inches(0.05), fill=GREEN_DARK)
        add_text(
            slide,
            gx + Inches(0.15), goals_y + Inches(0.12),
            goal_w - Inches(0.3), Inches(0.3),
            [title],
            font_size=13, bold=True, color=GREEN_DARK,
        )
        add_text(
            slide,
            gx + Inches(0.15), goals_y + Inches(0.42),
            goal_w - Inches(0.3), Inches(0.28),
            [subtitle],
            font_size=10, italic=True, color=GRAY_700,
        )
        blines = [{"text": "• " + b, "space_after": 3} for b in bullets]
        add_text(
            slide,
            gx + Inches(0.15), goals_y + Inches(0.75),
            goal_w - Inches(0.3), goal_h - Inches(0.85),
            blines,
            font_size=10, color=GRAY_800, line_spacing=1.2,
        )
        gx += goal_w + goal_gap

    # Bottom: 2 challenge cards
    cy = goals_y + goal_h + Inches(0.25)
    add_section_title(slide, CONTENT_PAD_X, cy, CONTENT_W, "核心挑战 (Challenges)")
    cy2 = cy + Inches(0.4)
    ch_h = Inches(1.0)
    ch_gap = Inches(0.2)
    ch_w = (CONTENT_W - ch_gap) / 2
    challenges = [
        (
            "C1 · 测试预言机缺失",
            "防御失效是 '静默' 逻辑漏洞 —— 无 crash、无报错; 传统 crash-oracle fuzzer 全盘失效。",
        ),
        (
            "C2 · 深层路径难触发",
            "GCC 千万行代码, 防御逻辑分散多 pass; 随机变异难以生成满足路径约束的 seed。",
        ),
    ]
    cx = CONTENT_PAD_X
    for title, desc in challenges:
        add_rect(
            slide, cx, cy2, ch_w, ch_h,
            fill=RED_LIGHT, line=RED_DARK, line_width=0.75,
        )
        add_rect(slide, cx, cy2, Inches(0.08), ch_h, fill=RED_DARK)
        add_text(
            slide,
            cx + Inches(0.2), cy2 + Inches(0.1),
            ch_w - Inches(0.3), Inches(0.32),
            [title],
            font_size=12, bold=True, color=RED_DARK,
        )
        add_text(
            slide,
            cx + Inches(0.2), cy2 + Inches(0.42),
            ch_w - Inches(0.3), ch_h - Inches(0.45),
            [desc],
            font_size=10, color=GRAY_800, line_spacing=1.3,
        )
        cx += ch_w + ch_gap

    # Footnote: idea origin
    foot_y = cy2 + ch_h + Inches(0.1)
    add_text(
        slide,
        CONTENT_PAD_X, foot_y, CONTENT_W, Inches(0.3),
        [
            {
                "runs": [
                    {"text": "思路起源: ", "bold": True, "color": GRAY_700, "size": 10},
                    {
                        "text": "CVE-2023-4039 (AArch64 canary 绕过) 启发我们系统化抽取防御机制的不变量。",
                        "color": GRAY_700,
                        "size": 10,
                    },
                ],
                "align": PP_ALIGN.LEFT,
            }
        ],
    )
    return slide


# ─────────────────────────────────────────────────────────────
# Slide 6 — 研究方案: 三步骤
# ─────────────────────────────────────────────────────────────
def slide_06_method(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[BLANK_LAYOUT])
    add_page_frame(slide, "5. 研究方案 — 三步骤工作流", page_no=5)

    y = CONTENT_TOP
    add_section_title(
        slide, CONTENT_PAD_X, y, CONTENT_W,
        "Invariants 调研 → Oracle 落地 → LLM + 覆盖率驱动 Fuzzing",
    )

    # 3 step cards
    sy = y + Inches(0.45)
    sh = Inches(2.35)
    arrow_w = Inches(0.45)
    card_w = (CONTENT_W - arrow_w * 2) / 3
    cx = CONTENT_PAD_X
    steps = [
        {
            "n": "1",
            "title": "Invariants 调研",
            "subtitle": "形式化捕获防御机制规范",
            "bullets": [
                "来源 A · 编译器源码注释 / pass 实现",
                "来源 B · GCC / LLVM / glibc 官方文档",
                "来源 C · 已确认历史 bug / CVE",
                "产物: 24 份跨机制 invariant survey",
            ],
        },
        {
            "n": "2",
            "title": "Oracle 落地",
            "subtitle": "invariant → 可执行 Checker",
            "bullets": [
                "InvariantChecker 接口: Enablement / Static / Dynamic",
                "MechanismOracle 分阶段聚合判定",
                "Polarity-aware: 负控 (-fno-*) 自动翻极性",
                "报告带 invariant ID + 源码引用",
            ],
        },
        {
            "n": "3",
            "title": "覆盖率 + LLM Fuzzing",
            "subtitle": "以防御机制源码覆盖率为信号",
            "bullets": [
                "CFG-guided target 选择 (插桩 GCC + .gcda)",
                "Mechanism contract + 分层 prompt 生成 seed",
                "uftrace 发散分析指导重试",
                "coverage / oracle / corpus 三方决策",
            ],
        },
    ]
    for i, s in enumerate(steps):
        add_rect(
            slide, cx, sy, card_w, sh,
            fill=GREEN_LIGHT, line=GREEN_DARK, line_width=0.75,
        )
        # step number circle
        circle_d = Inches(0.55)
        add_rect(
            slide, cx + Inches(0.15), sy + Inches(0.12),
            circle_d, circle_d,
            fill=GREEN_DARK, shape=MSO_SHAPE.OVAL,
        )
        add_text(
            slide,
            cx + Inches(0.15), sy + Inches(0.12),
            circle_d, circle_d,
            [s["n"]],
            font_size=18, bold=True, color=WHITE,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )
        add_text(
            slide,
            cx + Inches(0.8), sy + Inches(0.14),
            card_w - Inches(0.9), Inches(0.3),
            [s["title"]],
            font_size=13, bold=True, color=GREEN_DARK,
        )
        add_text(
            slide,
            cx + Inches(0.8), sy + Inches(0.42),
            card_w - Inches(0.9), Inches(0.25),
            [s["subtitle"]],
            font_size=9, italic=True, color=GRAY_700,
        )
        bullet_items = [{"text": "• " + b, "space_after": 3} for b in s["bullets"]]
        add_text(
            slide,
            cx + Inches(0.2), sy + Inches(0.8),
            card_w - Inches(0.35), sh - Inches(0.9),
            bullet_items,
            font_size=10, color=GRAY_800, line_spacing=1.25,
        )
        if i < 2:
            add_arrow(
                slide, cx + card_w, sy + sh / 2 - Inches(0.18),
                arrow_w, Inches(0.36),
                fill=GREEN_DARK,
            )
        cx += card_w + arrow_w

    # Bottom pipeline ribbon
    py = sy + sh + Inches(0.2)
    ph = Inches(0.5)
    add_rect(
        slide, CONTENT_PAD_X, py, CONTENT_W, ph,
        fill=YELLOW_LIGHT, line=YELLOW_ACCENT, line_width=0.75,
    )
    add_text(
        slide,
        CONTENT_PAD_X + Inches(0.15), py + Inches(0.04),
        CONTENT_W - Inches(0.3), Inches(0.24),
        ["Fuzzing 主循环"],
        font_size=11, bold=True, color=GREEN_DARK,
    )
    add_text(
        slide,
        CONTENT_PAD_X + Inches(0.15), py + Inches(0.26),
        CONTENT_W - Inches(0.3), ph - Inches(0.3),
        [
            "target → prompt → LLM → compile + gcov → execute (native / QEMU) → "
            "MechanismOracle → coverage / corpus / retry 反馈"
        ],
        font_size=10, color=GRAY_800,
    )
    return slide


# ─────────────────────────────────────────────────────────────
# Slide 7 — 技术贡献
# ─────────────────────────────────────────────────────────────
def slide_07_tech(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[BLANK_LAYOUT])
    add_page_frame(slide, "5.1 技术贡献", page_no=6)

    y = CONTENT_TOP
    add_section_title(slide, CONTENT_PAD_X, y, CONTENT_W, "支撑上述三步骤的三项原创工程贡献")

    cy = y + Inches(0.45)
    ch = CONTENT_H - Inches(0.6)
    gap = Inches(0.18)
    cw = (CONTENT_W - gap * 2) / 3

    tech_cards = [
        {
            "title": "T1 · 覆盖率插桩 GCC 工具链",
            "color_fill": GREEN_LIGHT,
            "color_border": GREEN_DARK,
            "color_title": GREEN_DARK,
            "bullets": [
                "多 ISA 交叉编译: x64 / aarch64 / loongarch64 / riscv64",
                "gcovr 读取 .gcda 构建防御机制源码覆盖率",
                "QEMU-user 运行交叉产物, 统一退出码协议",
                "与主循环异步: 每 seed 独立产物, 支持断点恢复",
            ],
        },
        {
            "title": "T2 · 多 invariant Oracle 框架",
            "color_fill": BLUE_LIGHT,
            "color_border": BLUE_DARK,
            "color_title": BLUE_DARK,
            "bullets": [
                "InvariantChecker 接口 (Enablement / Static / Dynamic)",
                "MechanismOracle 三阶段调度, 早退化早报错",
                "Polarity-aware: -fno-* 负控下自动翻极性",
                "Finding 带 invariant ID + 一手证据链接",
            ],
        },
        {
            "title": "T3 · Contract + 分层 Prompt",
            "color_fill": ORANGE_LIGHT,
            "color_border": ORANGE_DARK,
            "color_title": ORANGE_DARK,
            "bullets": [
                "三段式 prompt: base + understanding + contract",
                "RequiredMarkers 后置校验 LLM 响应",
                "Strategy ↔ Mechanism ↔ Oracle 绑定, 防错配",
                "FlagProfile 注入 CFLAGS 变种, 驱动 negative control",
            ],
        },
    ]
    cx = CONTENT_PAD_X
    for c in tech_cards:
        add_rect(
            slide, cx, cy, cw, ch,
            fill=c["color_fill"], line=c["color_border"], line_width=0.75,
        )
        add_rect(slide, cx, cy, cw, Inches(0.05), fill=c["color_border"])
        add_text(
            slide,
            cx + Inches(0.15), cy + Inches(0.15),
            cw - Inches(0.3), Inches(0.5),
            [c["title"]],
            font_size=13, bold=True, color=c["color_title"],
            line_spacing=1.1,
        )
        items = [{"text": "• " + b, "space_after": 5} for b in c["bullets"]]
        add_text(
            slide,
            cx + Inches(0.18), cy + Inches(0.8),
            cw - Inches(0.36), ch - Inches(0.9),
            items,
            font_size=10, color=GRAY_800, line_spacing=1.3,
        )
        cx += cw + gap
    return slide


# ─────────────────────────────────────────────────────────────
# Slide 8 — Oracle 实例 Canary
# ─────────────────────────────────────────────────────────────
def slide_08_canary_oracle(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[BLANK_LAYOUT])
    add_page_frame(slide, "6. Oracle 实例 — Stack Canary 二分检测", page_no=7)

    y = CONTENT_TOP
    add_section_title(slide, CONTENT_PAD_X, y, CONTENT_W, "4 个 Invariant Checker 覆盖 canary 机制的正确性契约")

    y2 = y + Inches(0.42)
    left_w = Inches(5.4)
    right_w = CONTENT_W - left_w - Inches(0.2)
    left_x = CONTENT_PAD_X
    right_x = left_x + left_w + Inches(0.2)

    # LEFT column: principle + 4 invariants list
    # Principle box
    pp_h = Inches(0.75)
    add_rect(slide, left_x, y2, left_w, pp_h, fill=GRAY_100, line=GRAY_300)
    add_text(
        slide,
        left_x + Inches(0.15), y2 + Inches(0.06),
        left_w - Inches(0.3), Inches(0.25),
        ["核心原理"], font_size=11, bold=True, color=GREEN_DARK,
    )
    add_text(
        slide,
        left_x + Inches(0.15), y2 + Inches(0.32),
        left_w - Inches(0.3), pp_h - Inches(0.34),
        [
            "渐进增加 buffer 填充 → 观察 ret / canary / buf 三者位置的单调退出状态变化 → 二分定位最小崩溃长度。"
        ],
        font_size=10, color=GRAY_800, line_spacing=1.25,
    )

    # Invariant table
    it_y = y2 + pp_h + Inches(0.12)
    it_h = CONTENT_H - pp_h - Inches(0.5)
    add_rect(slide, left_x, it_y, left_w, Inches(0.35), fill=GREEN_DARK)
    add_text(
        slide,
        left_x + Inches(0.15), it_y, left_w - Inches(0.3), Inches(0.35),
        ["已注册 InvariantChecker (4)"],
        font_size=11, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE,
    )
    inv_rows = [
        ("INV-SP-G01", "Static", "__stack_chk_fail 符号存在"),
        ("INV-SP-A01", "Static", "main 不应持有 canary 槽"),
        ("INV-SP-L01", "Dynamic", "buffer 溢出二分搜索 + 哨兵"),
        ("INV-SP-R03", "Dynamic", "epilogue 必须 scrub canary 寄存器 (新)"),
    ]
    row_h = (it_h - Inches(0.35)) / len(inv_rows)
    ry = it_y + Inches(0.35)
    for idx, (inv_id, kind, desc) in enumerate(inv_rows):
        bg = WHITE if idx % 2 == 0 else GREEN_LIGHT
        add_rect(slide, left_x, ry, left_w, row_h, fill=bg, line=GREEN_PALE, line_width=0.5)
        add_text(
            slide,
            left_x + Inches(0.12), ry, Inches(1.35), row_h,
            [inv_id], font_size=10, bold=True, color=GREEN_DARK, anchor=MSO_ANCHOR.MIDDLE,
        )
        add_text(
            slide,
            left_x + Inches(1.5), ry, Inches(0.9), row_h,
            [kind], font_size=10, italic=True, color=GRAY_700, anchor=MSO_ANCHOR.MIDDLE,
        )
        add_text(
            slide,
            left_x + Inches(2.4), ry,
            left_w - Inches(2.55), row_h,
            [desc], font_size=10, color=GRAY_800, anchor=MSO_ANCHOR.MIDDLE,
        )
        ry += row_h

    # RIGHT column: judgement table + sentinel
    # Judgement table
    jt_y = y2
    jt_h = Inches(1.75)
    add_rect(slide, right_x, jt_y, right_w, Inches(0.35), fill=GREEN_DARK)
    add_text(
        slide,
        right_x + Inches(0.1), jt_y, right_w - Inches(0.2), Inches(0.35),
        ["退出状态判定表 (INV-SP-L01)"],
        font_size=11, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE,
    )
    jrows = [
        ("有 SEED_RETURNED", "139", "Canary 绕过 (BUG)", RED_DARK),
        ("无", "139", "间接崩溃 (假阳性)", ORANGE_DARK),
        ("—", "134", "Canary 生效 (安全)", GREEN_DARK),
    ]
    header_widths = [Inches(1.35), Inches(0.75), right_w - Inches(2.1) - Inches(0.2)]
    # Column headers
    headers_y = jt_y + Inches(0.35)
    headers_h = Inches(0.3)
    add_rect(slide, right_x, headers_y, right_w, headers_h, fill=GREEN_PALE)
    hx = right_x + Inches(0.1)
    for label, w in zip(["stdout", "exit", "判定"], header_widths):
        add_text(
            slide, hx, headers_y, w, headers_h,
            [label], font_size=9, bold=True, color=GREEN_DARK,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )
        hx += w
    # Data rows
    jrh = (jt_h - Inches(0.35) - headers_h) / len(jrows)
    ry = headers_y + headers_h
    for idx, (stdout_v, exit_v, verdict, verdict_col) in enumerate(jrows):
        bg = WHITE if idx % 2 == 0 else GREEN_LIGHT
        add_rect(slide, right_x, ry, right_w, jrh, fill=bg, line=GREEN_PALE, line_width=0.5)
        hx = right_x + Inches(0.1)
        cells = [
            (stdout_v, GRAY_800, False, False),
            (exit_v, GRAY_800, True, False),
            (verdict, verdict_col, True, False),
        ]
        for (txt, col, b, it), w in zip(cells, header_widths):
            add_text(
                slide, hx, ry, w, jrh,
                [txt], font_size=10, bold=b, italic=it, color=col,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
            )
            hx += w
        ry += jrh

    # Sentinel / cross-ISA description
    sy = jt_y + jt_h + Inches(0.12)
    sh = CONTENT_H - jt_h - Inches(0.4)
    add_rect(slide, right_x, sy, right_w, sh, fill=YELLOW_LIGHT, line=YELLOW_ACCENT, line_width=0.5)
    add_text(
        slide,
        right_x + Inches(0.15), sy + Inches(0.08),
        right_w - Inches(0.3), Inches(0.3),
        ["哨兵 & 跨架构执行"],
        font_size=11, bold=True, color=GREEN_DARK,
    )
    add_text(
        slide,
        right_x + Inches(0.15), sy + Inches(0.36),
        right_w - Inches(0.3), sh - Inches(0.4),
        [
            {"text": "seed() 返回前打印 \"SEED_RETURNED\" —— 区分 '真 canary 绕过' 与 '函数内部间接崩溃'。", "space_after": 3},
            {"text": "0x41 填充在 x64 / aarch64 / loongarch64 / riscv64 均映射到未映射地址, 稳定触发 SIGSEGV。", "space_after": 3},
            {"text": "交叉编译 + QEMU-user 执行, 退出码统一由主循环判定。"},
        ],
        font_size=10, color=GRAY_800, line_spacing=1.25,
    )
    return slide


# ─────────────────────────────────────────────────────────────
# Slide 9 — 初步实验结果
# ─────────────────────────────────────────────────────────────
def slide_09_experiment(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[BLANK_LAYOUT])
    add_page_frame(slide, "7. 初步实验结果", page_no=8)

    y = CONTENT_TOP
    add_section_title(slide, CONTENT_PAD_X, y, CONTENT_W, "已在 GCC 上覆盖多防御机制 × 多 ISA, 发现 3 个真实漏洞")

    # Coverage matrix
    ty = y + Inches(0.42)
    th = Inches(1.6)
    header_h = Inches(0.32)
    # Columns: 机制 | x64 | aarch64 | loongarch64 | riscv64
    col_w = [Inches(2.6), Inches(1.42), Inches(1.42), Inches(1.52), Inches(1.42)]
    total_cw = sum(col_w, Inches(0))
    tx = CONTENT_PAD_X + (CONTENT_W - total_cw) / 2

    # Header row
    add_rect(slide, tx, ty, total_cw, header_h, fill=GREEN_DARK)
    hx = tx
    for lbl, w in zip(["防御机制", "x64", "aarch64", "loongarch64", "riscv64"], col_w):
        add_text(
            slide, hx, ty, w, header_h,
            [lbl], font_size=11, bold=True, color=WHITE,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )
        hx += w
    # Rows
    rows = [
        ("Stack Canary", ["✓", "✓", "✓", "✓"]),
        ("FORTIFY_SOURCE", ["✓", "—", "—", "—"]),
        ("CET-IBT", ["✓", "N/A", "N/A", "N/A"]),
    ]
    row_h = (th - header_h) / len(rows)
    ry = ty + header_h
    for idx, (name, marks) in enumerate(rows):
        bg = WHITE if idx % 2 == 0 else GREEN_LIGHT
        add_rect(slide, tx, ry, total_cw, row_h, fill=bg, line=GREEN_PALE, line_width=0.5)
        hx = tx
        add_text(
            slide, hx + Inches(0.15), ry, col_w[0] - Inches(0.3), row_h,
            [name], font_size=11, bold=True, color=GREEN_DARK,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        hx += col_w[0]
        for m, w in zip(marks, col_w[1:]):
            if m == "✓":
                col = GREEN_DARK
                bold = True
            elif m == "N/A":
                col = GRAY_500
                bold = False
            else:
                col = GRAY_500
                bold = False
            add_text(
                slide, hx, ry, w, row_h,
                [m], font_size=13, bold=bold, color=col,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
            )
            hx += w
        ry += row_h

    # Bugs found
    by = ty + th + Inches(0.2)
    add_section_title(
        slide, CONTENT_PAD_X, by, CONTENT_W,
        "已发现 3 个 GCC 防御实现漏洞 (均已上报, 其中 1 个已进 master)",
        font_size=13,
    )
    bugs_y = by + Inches(0.4)
    bugs = [
        (
            "Bug 1",
            "Canary epilogue 残留",
            "长尾 ISA (loongarch64 / mips / xtensa / csky) 的 generic fallback "
            "未在 ret 前 scrub caller-saved 寄存器, guard 值泄露。",
        ),
        (
            "Bug 2",
            "FORTIFY bound 未更新",
            "tree-ssa-strlen 把 __strcat_chk 改写为 __stpcpy_chk 时, "
            "arg3 (objsize) 仍按原缓冲区大小传递, 保护静默退化。",
        ),
        (
            "Bug 3",
            "CET-IBT 谓词不充分",
            "movabsq 立即数若在非最低字节位置包含 ENDBR64 字节序列, "
            "可构造出编译器未预期的 landing pad, 削弱 IBT。",
        ),
    ]
    bug_h = Inches(0.95)
    bug_gap = Inches(0.1)
    bw = (CONTENT_W - bug_gap * 2) / 3
    bx = CONTENT_PAD_X
    for tag, title, desc in bugs:
        add_rect(
            slide, bx, bugs_y, bw, bug_h,
            fill=RED_LIGHT, line=RED_DARK, line_width=0.75,
        )
        add_rect(slide, bx, bugs_y, bw, Inches(0.05), fill=RED_DARK)
        add_text(
            slide,
            bx + Inches(0.12), bugs_y + Inches(0.1),
            bw - Inches(0.2), Inches(0.28),
            [
                {
                    "runs": [
                        {"text": tag + " · ", "bold": True, "color": RED_DARK, "size": 11},
                        {"text": title, "bold": True, "color": GRAY_800, "size": 11},
                    ]
                }
            ],
        )
        add_text(
            slide,
            bx + Inches(0.12), bugs_y + Inches(0.42),
            bw - Inches(0.2), bug_h - Inches(0.46),
            [desc],
            font_size=9, color=GRAY_800, line_spacing=1.3,
        )
        bx += bw + bug_gap

    # Future work strip
    fw_y = bugs_y + bug_h + Inches(0.15)
    fw_h = CONTENT_H - (fw_y - CONTENT_TOP) - Inches(0.05)
    add_rect(slide, CONTENT_PAD_X, fw_y, CONTENT_W, fw_h,
             fill=YELLOW_LIGHT, line=YELLOW_ACCENT, line_width=0.75)
    add_rect(slide, CONTENT_PAD_X, fw_y, Inches(0.08), fw_h, fill=YELLOW_ACCENT)
    add_text(
        slide,
        CONTENT_PAD_X + Inches(0.2), fw_y + Inches(0.04),
        CONTENT_W - Inches(0.3), Inches(0.26),
        ["后续实验"],
        font_size=11, bold=True, color=GREEN_DARK,
    )
    add_text(
        slide,
        CONTENT_PAD_X + Inches(0.2), fw_y + Inches(0.3),
        CONTENT_W - Inches(0.3), fw_h - Inches(0.34),
        [
            "扩展至 LLVM 工具链 + 补全 GCC 其他防御机制 (SafeStack · CFI · PAC · BTI · ShadowStack · StackClash 等)。"
        ],
        font_size=10, color=GRAY_800, line_spacing=1.25,
    )
    return slide


# ─────────────────────────────────────────────────────────────
# Slide 10 — 预期成果
# ─────────────────────────────────────────────────────────────
def slide_10_expected(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[BLANK_LAYOUT])
    add_page_frame(slide, "8. 预期成果", page_no=9)

    y = CONTENT_TOP
    add_section_title(slide, CONTENT_PAD_X, y, CONTENT_W, "学术产出 + 工程产出 + 社区影响")

    cy = y + Inches(0.5)
    ch = CONTENT_H - Inches(0.6)
    gap = Inches(0.2)
    cw = (CONTENT_W - gap) / 2
    cols = [
        {
            "title": "学术产出",
            "color": GREEN_DARK,
            "bg": GREEN_LIGHT,
            "bullets": [
                "1 篇 CCF-A / B 级会议或期刊论文 (目标方向: 安全 / 系统)",
                "开源 DeFuzz 工具 (Go + Python, 多 ISA 支持)",
                "LLM-driven 编译器防御机制测试方法论",
                "跨机制 invariant survey 作为独立贡献 (24 份)",
            ],
        },
        {
            "title": "工程 & 社区",
            "color": ORANGE_DARK,
            "bg": ORANGE_LIGHT,
            "bullets": [
                "向 GCC / LLVM 上游报告真实防御实现漏洞",
                "支持 GCC + LLVM 双工具链, 便于回归",
                "可扩展的 Invariant Oracle 体系, 覆盖面持续增长",
                "自动化流水线 (插桩编译 / QEMU 执行 / 报告)",
            ],
        },
    ]
    cx = CONTENT_PAD_X
    for c in cols:
        add_rect(slide, cx, cy, cw, ch, fill=c["bg"], line=c["color"], line_width=0.75)
        add_rect(slide, cx, cy, cw, Inches(0.05), fill=c["color"])
        add_text(
            slide,
            cx + Inches(0.2), cy + Inches(0.18),
            cw - Inches(0.4), Inches(0.4),
            [c["title"]],
            font_size=16, bold=True, color=c["color"],
        )
        items = [{"text": "• " + b, "space_after": 7} for b in c["bullets"]]
        add_text(
            slide,
            cx + Inches(0.25), cy + Inches(0.75),
            cw - Inches(0.45), ch - Inches(0.85),
            items,
            font_size=11, color=GRAY_800, line_spacing=1.35,
        )
        cx += cw + gap
    return slide


# ─────────────────────────────────────────────────────────────
# Slide 11 — 进度安排
# ─────────────────────────────────────────────────────────────
def slide_11_schedule(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[BLANK_LAYOUT])
    add_page_frame(slide, "9. 进度安排", page_no=10)

    y = CONTENT_TOP
    add_section_title(slide, CONTENT_PAD_X, y, CONTENT_W, "当前节点 → 毕业答辩的时间线")

    ty = y + Inches(0.45)
    th = CONTENT_H - Inches(0.55)
    add_rect(slide, CONTENT_PAD_X, ty, CONTENT_W, th, fill=GRAY_100, line=GRAY_300)

    # Vertical timeline bar
    bar_x = CONTENT_PAD_X + Inches(1.9)
    add_rect(slide, bar_x, ty + Inches(0.1), Inches(0.06), th - Inches(0.2), fill=GREEN_DARK)

    entries = [
        {
            "period": "2025.11–2026.04",
            "label": "已完成",
            "color": GREEN_DARK,
            "tasks": [
                "Canary oracle 4 个 checker · Mechanism contract 架构 · 3 个 GCC bug 上报",
            ],
            "done": True,
        },
        {
            "period": "2026.05",
            "label": "LLVM 接入",
            "color": GREEN_MID,
            "tasks": ["接入 Clang/LLVM 工具链 + 扩展 FORTIFY / CFI oracle"],
            "done": False,
        },
        {
            "period": "2026.06",
            "label": "大规模实验",
            "color": GREEN_MID,
            "tasks": ["多机制 × 多 ISA × 多版本完整实验 + 数据收集"],
            "done": False,
        },
        {
            "period": "2026.07",
            "label": "论文初稿",
            "color": ORANGE_DARK,
            "tasks": ["撰写论文初稿 + invariant survey 整理"],
            "done": False,
        },
        {
            "period": "2026.08",
            "label": "投稿 & 修订",
            "color": ORANGE_DARK,
            "tasks": ["会议投稿, 内部 review 迭代 + 补充实验"],
            "done": False,
        },
        {
            "period": "2026.09–10",
            "label": "毕业论文",
            "color": RED_DARK,
            "tasks": ["毕业论文撰写 + 答辩准备"],
            "done": False,
        },
    ]
    n = len(entries)
    avail_h = th - Inches(0.3)
    row_h = avail_h / n
    ry = ty + Inches(0.15)
    for e in entries:
        # Dot on the timeline
        dot_d = Inches(0.2)
        add_rect(
            slide,
            bar_x + Inches(0.03) - dot_d / 2,
            ry + row_h / 2 - dot_d / 2,
            dot_d, dot_d,
            fill=e["color"], shape=MSO_SHAPE.OVAL,
        )
        # Period (left of bar)
        add_text(
            slide,
            CONTENT_PAD_X + Inches(0.2), ry, Inches(1.65), row_h,
            [e["period"]],
            font_size=11, bold=True, color=e["color"], align=PP_ALIGN.RIGHT,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        # Label + tasks (right of bar)
        text_x = bar_x + Inches(0.25)
        text_w = CONTENT_W - (text_x - CONTENT_PAD_X) - Inches(0.15)
        add_text(
            slide,
            text_x, ry + Inches(0.04), text_w, Inches(0.26),
            [
                {
                    "runs": [
                        {"text": e["label"], "bold": True, "color": e["color"], "size": 12},
                        {"text": "  ", "size": 12},
                        {
                            "text": "✓ 已完成" if e["done"] else "▷ 计划",
                            "italic": True,
                            "color": GRAY_600,
                            "size": 9,
                        },
                    ]
                }
            ],
        )
        tasks_str = " · ".join(e["tasks"])
        add_text(
            slide,
            text_x, ry + Inches(0.32), text_w, row_h - Inches(0.36),
            [tasks_str],
            font_size=10, color=GRAY_800, line_spacing=1.2,
        )
        ry += row_h
    return slide


# ─────────────────────────────────────────────────────────────
# Slide 12 — Thanks
# ─────────────────────────────────────────────────────────────
def slide_12_thanks(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[BLANK_LAYOUT])
    # Accent bars
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.12), fill=GREEN_DARK)
    add_rect(slide, 0, SLIDE_H - Inches(0.08), SLIDE_W, Inches(0.08), fill=YELLOW_ACCENT)

    logo_w = Inches(0.75)
    slide.shapes.add_picture(
        str(LOGO_PATH),
        (SLIDE_W - logo_w) / 2,
        Inches(1.15),
        width=logo_w,
        height=logo_w,
    )
    add_text(
        slide,
        Inches(0.5), Inches(2.35), Inches(9), Inches(1.0),
        ["感谢各位老师的聆听!"],
        font_size=44, bold=True, color=GREEN_DARK,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        Inches(0.5), Inches(3.65), Inches(9), Inches(0.5),
        ["敬请批评指正"],
        font_size=20, italic=True, color=GRAY_600,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )
    return slide


ALL_SLIDES = [
    slide_01_title,
    slide_02_background,
    slide_03_status,
    slide_04_motivation,
    slide_05_goals,
    slide_06_method,
    slide_07_tech,
    slide_08_canary_oracle,
    slide_09_experiment,
    slide_10_expected,
    slide_11_schedule,
    slide_12_thanks,
]
