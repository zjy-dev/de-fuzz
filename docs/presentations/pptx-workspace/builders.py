"""Per-slide builders for the DeFuzz 开题答辩 deck.

章节结构:
    01 研究背景与意义   slides 1-3
    02 研究现状         slide 4
    03 研究目标与内容   slides 5-7
    04 研究方案         slides 8-10
    05 现有成果         slides 11-13

页码由 theme.PageCounter 自动维护; 扉页 builder 不调 pc.next(),
内容页调 pc.next() 取页号传给 add_page_frame.
"""
from __future__ import annotations

from pathlib import Path

from pptx.util import Inches, Pt

from theme import (
    BLUE_DARK,
    BLUE_LIGHT,
    CONTENT_H,
    CONTENT_PAD_X,
    CONTENT_PAD_Y,
    CONTENT_TOP,
    CONTENT_W,
    FOOTER_H,
    GRAY_100,
    GRAY_300,
    GRAY_500,
    GRAY_600,
    GRAY_700,
    GRAY_800,
    GREEN_DARK,
    GREEN_LIGHT,
    GREEN_PALE,
    HEADER_H,
    LOGO_PATH,
    MSO_ANCHOR,
    MSO_SHAPE,
    PP_ALIGN,
    RED_DARK,
    RED_LIGHT,
    SLIDE_H,
    SLIDE_W,
    WHITE,
    YELLOW_ACCENT,
    YELLOW_LIGHT,
    add_arrow,
    add_page_frame,
    add_rect,
    add_section_title,
    add_text,
)

BLANK_LAYOUT = 6  # python-pptx default blank layout index

ASSETS_DIR = Path(__file__).resolve().parent / "assets"
IMG_CVE = ASSETS_DIR / "cve_growth.zh.png"
IMG_STACK = ASSETS_DIR / "stack-layout.zh.png"
IMG_MATRIX = ASSETS_DIR / "defense-matrix.zh.png"
IMG_PIPELINE = ASSETS_DIR / "pipeline.zh.png"
IMG_ARCH = ASSETS_DIR / "architecture.zh.png"
IMG_ORACLE = ASSETS_DIR / "oracle-flow.zh.png"
IMG_GCC_DEV = ASSETS_DIR / "gcc-dev-severe.png"


# ─────────────────────────────────────────────────────────────
# 通用: 章节扉页渲染
# ─────────────────────────────────────────────────────────────
def _render_section_divider(slide, number: str, title_zh: str, title_en: str, footnote: str) -> None:
    """白底 + 绿/黄装饰条 + 大号章节号 + 主副标题 + 底部脚注 (5 章共用)."""
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.12), fill=GREEN_DARK)
    add_rect(slide, 0, SLIDE_H - Inches(0.08), SLIDE_W, Inches(0.08), fill=YELLOW_ACCENT)

    # 大号章节号
    add_text(
        slide,
        Inches(0.7), Inches(1.55), Inches(2.2), Inches(1.6),
        [number],
        font_size=110, bold=True, color=GREEN_PALE,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
    )
    # 黄色短线分隔
    add_rect(
        slide,
        Inches(3.0), Inches(2.32), Inches(0.08), Inches(0.95),
        fill=YELLOW_ACCENT,
    )
    # 主标题
    add_text(
        slide,
        Inches(3.3), Inches(2.25), Inches(6.2), Inches(0.7),
        [title_zh],
        font_size=40, bold=True, color=GREEN_DARK,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
    )
    # 英文副标题
    add_text(
        slide,
        Inches(3.3), Inches(2.95), Inches(6.2), Inches(0.4),
        [title_en],
        font_size=18, italic=True, color=GRAY_600,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
    )
    # 底部脚注
    add_text(
        slide,
        Inches(0.7), SLIDE_H - Inches(0.7),
        Inches(8.6), Inches(0.32),
        [footnote],
        font_size=12, italic=True, color=GRAY_700,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
    )


def _new_slide(pres):
    return pres.slides.add_slide(pres.slide_layouts[BLANK_LAYOUT])


# ═════════════════════════════════════════════════════════════
# 章节 01 — 研究背景与意义 (slides 1-3)
# ═════════════════════════════════════════════════════════════


def slide_section_01_divider(pres, pc):
    slide = _new_slide(pres)
    _render_section_divider(
        slide,
        number="01",
        title_zh="研究背景与意义",
        title_en="Background & Motivation",
        footnote="漏洞过载环境下，编译器防御机制如何成为系统安全的坚实后盾？",
    )
    return slide


def slide_01_cve_growth(pres, pc):
    slide = _new_slide(pres)
    add_page_frame(slide, "研究背景与意义 — 漏洞过载与防御底座", page_no=pc.next())

    y = CONTENT_TOP
    add_section_title(
        slide, CONTENT_PAD_X, y, CONTENT_W,
        "研究背景：漏洞过载环境下，防御机制愈发重要",
    )

    body_top = y + Inches(0.42)
    body_h = CONTENT_H - Inches(0.5) - Inches(0.3)

    chart_w = Inches(5.85)
    chart_h = Inches(5.85 / 2.295)
    chart_x = CONTENT_PAD_X
    chart_y = body_top + (body_h - chart_h) / 2
    slide.shapes.add_picture(str(IMG_CVE), chart_x, chart_y, width=chart_w, height=chart_h)
    add_text(
        slide,
        chart_x, chart_y - Inches(0.24),
        chart_w, Inches(0.22),
        ["图 1: 近 10 年 CVE 年度发布数"],
        font_size=10, italic=True, color=GRAY_600,
        align=PP_ALIGN.CENTER,
    )

    right_x = CONTENT_PAD_X + chart_w + Inches(0.22)
    right_w = CONTENT_W - chart_w - Inches(0.22)
    card_h = (body_h - Inches(0.18)) / 2

    ay = body_top
    add_rect(slide, right_x, ay, right_w, card_h, fill=GREEN_LIGHT, line=GREEN_DARK, line_width=0.75)
    add_rect(slide, right_x, ay, right_w, Inches(0.05), fill=GREEN_DARK)
    add_text(
        slide,
        right_x + Inches(0.15), ay + Inches(0.12),
        right_w - Inches(0.3), Inches(0.3),
        ["防御机制：重要性愈发显著"],
        font_size=12, bold=True, color=GREEN_DARK,
    )
    add_text(
        slide,
        right_x + Inches(0.15), ay + Inches(0.5),
        right_w - Inches(0.3), card_h - Inches(0.55),
        [
            {
                "runs": [
                    {"text": "防御机制（栈金丝雀、CFI 等）是软件安全的 \"兜底层\"——在漏洞", "size": 10.5},
                    {"text": "被触发时", "size": 10.5, "bold": True, "color": GREEN_DARK},
                    {"text": ", 仍能挡住攻击者拿到控制权。", "size": 10.5},
                ]
            },
        ],
        font_size=10.5, color=GRAY_800, line_spacing=1.3,
    )

    by = ay + card_h + Inches(0.18)
    add_rect(slide, right_x, by, right_w, card_h, fill=YELLOW_LIGHT, line=YELLOW_ACCENT, line_width=0.75)
    add_rect(slide, right_x, by, right_w, Inches(0.05), fill=YELLOW_ACCENT)
    add_text(
        slide,
        right_x + Inches(0.15), by + Inches(0.12),
        right_w - Inches(0.3), Inches(0.3),
        ["漏洞难以根除"],
        font_size=12, bold=True, color=GREEN_DARK,
    )
    add_text(
        slide,
        right_x + Inches(0.15), by + Inches(0.5),
        right_w - Inches(0.3), card_h - Inches(0.55),
        [
            {
                "runs": [
                    {"text": "CVE 数量逐年攀升, Google 等大厂已经意识到代码 Bug ", "size": 10.5},
                    {"text": "不可能被彻底修干净", "size": 10.5, "bold": True, "color": RED_DARK},
                    {"text": "。", "size": 10.5},
                ]
            }
        ],
        font_size=10.5, color=GRAY_800, line_spacing=1.35,
    )

    src_y = SLIDE_H - FOOTER_H - Inches(0.28)
    add_text(
        slide,
        CONTENT_PAD_X, src_y, CONTENT_W, Inches(0.22),
        ["数据来源: NVD · cve.org · JerryGamblin 2024–2025 CVE Data Review"],
        font_size=8.5, italic=True, color=GRAY_500,
        align=PP_ALIGN.LEFT,
    )
    return slide


def slide_02_canary_aarch64(pres, pc):
    slide = _new_slide(pres)
    add_page_frame(slide, "研究背景与意义 — 防御机制本身也会出错", page_no=pc.next())

    y = CONTENT_TOP
    add_section_title(
        slide, CONTENT_PAD_X, y, CONTENT_W,
        "现实困境：防御机制本身的脆弱性——栈金丝雀（Canary） 在 AArch64 上失效",
    )

    body_top = y + Inches(0.42)
    body_h = CONTENT_H - Inches(0.5) - Inches(0.28)

    img_h = body_h
    img_w = int(body_h * 0.804)
    img_x = CONTENT_PAD_X + Inches(0.1)
    img_y = body_top
    slide.shapes.add_picture(str(IMG_STACK), img_x, img_y, width=img_w, height=img_h)
    add_text(
        slide,
        img_x, body_top + body_h + Inches(0.04),
        img_w, Inches(0.22),
        ["图 2: AArch64 栈帧布局 (CVE-2023-4039)"],
        font_size=9.5, italic=True, color=GRAY_600,
        align=PP_ALIGN.CENTER,
    )

    right_x = img_x + img_w + Inches(0.25)
    right_w = CONTENT_W + CONTENT_PAD_X - right_x
    block_gap = Inches(0.18)
    block_h = (body_h - block_gap) / 2

    ay = body_top
    add_rect(slide, right_x, ay, right_w, block_h, fill=GREEN_LIGHT, line=GREEN_DARK, line_width=0.75)
    add_rect(slide, right_x, ay, Inches(0.08), block_h, fill=GREEN_DARK)
    add_text(
        slide,
        right_x + Inches(0.2), ay + Inches(0.1),
        right_w - Inches(0.3), Inches(0.3),
        ["Canary 介绍"],
        font_size=12, bold=True, color=GREEN_DARK,
    )
    add_text(
        slide,
        right_x + Inches(0.2), ay + Inches(0.46),
        right_w - Inches(0.3), block_h - Inches(0.5),
        [
            {
                "runs": [
                    {"text": "编译器在栈帧的", "size": 10.5},
                    {"text": "缓冲区与返回地址之间", "size": 10.5, "bold": True, "color": GREEN_DARK},
                    {"text": "放置随机值哨兵（金丝雀），以此来保护函数返回地址，阻止栈溢出（Buffer Overflow）攻击。", "size": 10.5},
                ],
                "space_after": 4,
            },
            {"text": "函数返回前检查金丝雀是否被改写, 一旦改写就报错强行终止程序。", "space_after": 4},
            {
                "runs": [
                    {"text": "主流包管理器在构建官方软件包时，均已", "size": 10.5},
                    {"text": "默认启用", "size": 10.5, "bold": True, "color": GREEN_DARK},
                    {"text": "栈金丝雀。", "size": 10.5},
                ]
            },
        ],
        color=GRAY_800, line_spacing=1.35,
    )

    by = ay + block_h + block_gap
    add_rect(slide, right_x, by, right_w, block_h, fill=RED_LIGHT, line=RED_DARK, line_width=0.75)
    add_rect(slide, right_x, by, Inches(0.08), block_h, fill=RED_DARK)
    add_text(
        slide,
        right_x + Inches(0.2), by + Inches(0.1),
        right_w - Inches(0.3), Inches(0.3),
        ["AArch64 上失效的根因 (CVE-2023-4039)"],
        font_size=12, bold=True, color=RED_DARK,
    )
    add_text(
        slide,
        right_x + Inches(0.2), by + Inches(0.46),
        right_w - Inches(0.3), block_h - Inches(0.5),
        [
            {
                "runs": [
                    {"text": "GCC 固定把 Canary 放在", "size": 10.5},
                    {"text": "局部变量区的最上方", "size": 10.5, "bold": True, "color": RED_DARK},
                    {"text": "。", "size": 10.5},
                ],
                "space_after": 4,
            },
            {
                "runs": [
                    {"text": "但 AArch64 栈帧把变长缓冲区（ VLA） 放在", "size": 10.5},
                    {"text": "返回地址之下", "size": 10.5, "bold": True, "color": RED_DARK},
                    {"text": ", ", "size": 10.5},
                    {"text": "金丝雀", "size": 10.5, "bold": True, "color": RED_DARK},
                    {"text": "被", "size": 10.5},
                    {"text": "绕开", "size": 10.5, "bold": True, "color": RED_DARK},
                    {"text": "。", "size": 10.5},
                ],
                "space_after": 4,
            },
            {
                "runs": [
                    {"text": "也就是说，防御机制的代码实现会与指令集架构（ISA）的设计出现冲突从而", "size": 10.5},
                    {"text": "失效", "size": 10.5, "bold": True, "color": RED_DARK},
                    {"text": "。", "size": 10.5},
                ],
            },
        ],
        color=GRAY_800, line_spacing=1.35,
    )

    src_y = SLIDE_H - FOOTER_H - Inches(0.28)
    add_text(
        slide,
        CONTENT_PAD_X, src_y, CONTENT_W, Inches(0.22),
        [
            "这种失效具有破坏性的“双重无声”特征：编译不报错，运行也不异常。直到 2023 年才被外部发现。"
        ],
        font_size=8.5, italic=True, color=GRAY_500,
        align=PP_ALIGN.LEFT,
    )
    return slide


def slide_03_core_challenge(pres, pc):
    slide = _new_slide(pres)
    add_page_frame(slide, "研究背景与意义 — 核心挑战", page_no=pc.next())

    y = CONTENT_TOP
    add_section_title(
        slide, CONTENT_PAD_X, y, CONTENT_W,
        "核心挑战：为何现有的检测工具行不通？",
    )

    body_top = y + Inches(0.42)
    ribbon_h = Inches(0.55)
    body_h = CONTENT_H - Inches(0.5) - ribbon_h - Inches(0.18)

    img_w = Inches(5.0)
    img_h = Inches(5.0 / 1.988)
    img_x = CONTENT_PAD_X
    img_y = body_top + (body_h - img_h - Inches(0.26)) / 2
    slide.shapes.add_picture(str(IMG_MATRIX), img_x, img_y, width=img_w, height=img_h)
    add_text(
        slide,
        img_x, img_y + img_h + Inches(0.04),
        img_w, Inches(0.22),
        ["图 3: 防御机制 × ISA 矩阵与漏洞上报情况"],
        font_size=9.5, italic=True, color=GRAY_600,
        align=PP_ALIGN.CENTER,
    )

    right_x = CONTENT_PAD_X + img_w + Inches(0.25)
    right_w = CONTENT_W + CONTENT_PAD_X - right_x
    card_gap = Inches(0.18)
    card_h = (body_h - card_gap) / 2

    ay = body_top
    add_rect(slide, right_x, ay, right_w, card_h, fill=GREEN_LIGHT, line=GREEN_DARK, line_width=0.75)
    add_rect(slide, right_x, ay, Inches(0.08), card_h, fill=GREEN_DARK)
    add_text(
        slide,
        right_x + Inches(0.2), ay + Inches(0.1),
        right_w - Inches(0.3), Inches(0.3),
        ["问题空间巨大：机制与架构的组合爆炸"],
        font_size=12, bold=True, color=GREEN_DARK,
    )
    add_text(
        slide,
        right_x + Inches(0.2), ay + Inches(0.46),
        right_w - Inches(0.3), card_h - Inches(0.5),
        [
            {"text": "防御机制和 ISA 种类多，且越来越多。", "space_after": 4},
            {"text": "“防御机制 × ISA”的二维空间十分庞大，且结合起来逻辑复杂，人工审计难以覆盖。"},
        ],
        font_size=10.5, color=GRAY_800, line_spacing=1.35,
    )

    by = ay + card_h + card_gap
    add_rect(slide, right_x, by, right_w, card_h, fill=RED_LIGHT, line=RED_DARK, line_width=0.75)
    add_rect(slide, right_x, by, Inches(0.08), card_h, fill=RED_DARK)
    add_text(
        slide,
        right_x + Inches(0.2), by + Inches(0.1),
        right_w - Inches(0.3), Inches(0.3),
        ["判定空白 (Oracle 难题)：如何察觉静默失效？"],
        font_size=12, bold=True, color=RED_DARK,
    )
    add_text(
        slide,
        right_x + Inches(0.2), by + Inches(0.46),
        right_w - Inches(0.3), card_h - Inches(0.5),
        [
            {
                "runs": [
                    {"text": "程序", "size": 10.5},
                    {"text": "不崩溃", "size": 10.5, "bold": True, "color": RED_DARK},
                    {"text": "，功能也", "size": 10.5},
                    {"text": "正常", "size": 10.5, "bold": True, "color": RED_DARK},
                    {"text": "，但安全契约已经", "size": 10.5},
                    {"text": "失守", "size": 10.5, "bold": True, "color": RED_DARK},
                    {"text": "。传统的漏洞挖掘方法对此", "size": 10.5},
                    {"text": "毫无察觉", "size": 10.5, "bold": True, "color": RED_DARK},
                    {"text": "。", "size": 10.5},
                ]
            }
        ],
        font_size=10.5, color=GRAY_800, line_spacing=1.35,
    )

    rb_y = SLIDE_H - FOOTER_H - ribbon_h - Inches(0.04)
    add_rect(slide, CONTENT_PAD_X, rb_y, CONTENT_W, ribbon_h, fill=YELLOW_LIGHT, line=YELLOW_ACCENT, line_width=0.75)
    add_rect(slide, CONTENT_PAD_X, rb_y, Inches(0.08), ribbon_h, fill=YELLOW_ACCENT)
    add_text(
        slide,
        CONTENT_PAD_X + Inches(0.2), rb_y + Inches(0.05),
        CONTENT_W - Inches(0.3), Inches(0.26),
        ["注"],
        font_size=11, bold=True, color=GREEN_DARK,
    )
    add_text(
        slide,
        CONTENT_PAD_X + Inches(0.2), rb_y + Inches(0.3),
        CONTENT_W - Inches(0.3), ribbon_h - Inches(0.32),
        ["深绿色的格点均存在通过本课题挖掘，且已上报且被 GCC 官方确认/修复的漏洞。"],
        font_size=10, color=GRAY_800, line_spacing=1.3,
    )
    return slide


# ═════════════════════════════════════════════════════════════
# 章节 02 — 研究现状 (slide 4)
# ═════════════════════════════════════════════════════════════


def slide_section_02_divider(pres, pc):
    slide = _new_slide(pres)
    _render_section_divider(
        slide,
        number="02",
        title_zh="研究现状",
        title_en="Related Work",
        footnote="静态分析与现有 Fuzz 框架为何在防御机制的静默失效面前集体失灵？",
    )
    return slide


def slide_04_status_methods(pres, pc):
    slide = _new_slide(pres)
    add_page_frame(slide, "研究现状 — 现有检测手段的局限", page_no=pc.next())

    y = CONTENT_TOP
    add_section_title(
        slide, CONTENT_PAD_X, y, CONTENT_W,
        "研究现状：现有自动化方法为何集体失灵？",
    )

    body_top = y + Inches(0.42)
    ribbon_h = Inches(0.65)
    body_h = CONTENT_H - Inches(0.5) - ribbon_h - Inches(0.18)

    # 左栏: 静态分析 (红框)
    left_w = Inches(3.4)
    left_x = CONTENT_PAD_X
    add_rect(slide, left_x, body_top, left_w, body_h, fill=RED_LIGHT, line=RED_DARK, line_width=0.75)
    add_rect(slide, left_x, body_top, left_w, Inches(0.05), fill=RED_DARK)
    add_text(
        slide,
        left_x + Inches(0.18), body_top + Inches(0.15),
        left_w - Inches(0.36), Inches(0.32),
        ["静态分析：难以应对的多样性"],
        font_size=13, bold=True, color=RED_DARK,
    )
    add_text(
        slide,
        left_x + Inches(0.18), body_top + Inches(0.55),
        left_w - Inches(0.36), body_h - Inches(0.65),
        [
            {
                "runs": [
                    {"text": "传统静态分析通过确定的代码", "size": 10.5},
                    {"text": "模式", "size": 10.5, "bold": True, "color": RED_DARK},
                    {"text": "去搜索海量的代码空间，从而发现同类漏洞。", "size": 10.5},
                ],
                "space_after": 6,
            },
            {
                "runs": [
                    {"text": "但防御机制失效的形态千奇百怪，大多数案例都有自己", "size": 10.5},
                    {"text": "独特的模式", "size": 10.5, "bold": True, "color": RED_DARK},
                    {"text": "，静态分析难以进行。", "size": 10.5},
                ],
                "space_after": 6,
            },
            {
                "runs": [
                    {"text": "因此本课题采用动态", "size": 10.5},
                    {"text": "模糊测试", "size": 10.5, "bold": True, "color": GREEN_DARK},
                    {"text": "（Fuzz）的方案。", "size": 10.5},
                ],
            },
        ],
        font_size=10.5, color=GRAY_800, line_spacing=1.4,
    )

    # 右栏: LLM 驱动的 fuzz (上下两个绿框)
    right_x = left_x + left_w + Inches(0.22)
    right_w = CONTENT_W + CONTENT_PAD_X - right_x
    card_gap = Inches(0.15)
    card_h = (body_h - card_gap) / 2

    # 卡 A: WhiteFox
    add_rect(slide, right_x, body_top, right_w, card_h, fill=GREEN_LIGHT, line=GREEN_DARK, line_width=0.75)
    add_rect(slide, right_x, body_top, Inches(0.08), card_h, fill=GREEN_DARK)
    add_text(
        slide,
        right_x + Inches(0.2), body_top + Inches(0.1),
        right_w - Inches(0.3), Inches(0.3),
        [
            {
                "runs": [
                    {"text": "WhiteFox", "bold": True, "color": GREEN_DARK, "size": 13},
                    {"text": "  · OSDI'24 ", "italic": True, "color": GRAY_600, "size": 10},
                ]
            }
        ],
    )
    add_text(
        slide,
        right_x + Inches(0.2), body_top + Inches(0.46),
        right_w - Inches(0.3), card_h - Inches(0.5),
        [
            {"text": "种子生成：用 LLM 归纳函数的代码模式，根据模式生成种子。", "space_after": 4},
            {"text": "预言机：常规的编译器崩溃或误编译，并不关心安全契约。"},
        ],
        font_size=10, color=GRAY_800, line_spacing=1.3,
    )

    # 卡 B: HLPFuzz
    by = body_top + card_h + card_gap
    add_rect(slide, right_x, by, right_w, card_h, fill=GREEN_LIGHT, line=GREEN_DARK, line_width=0.75)
    add_rect(slide, right_x, by, Inches(0.08), card_h, fill=GREEN_DARK)
    add_text(
        slide,
        right_x + Inches(0.2), by + Inches(0.1),
        right_w - Inches(0.3), Inches(0.3),
        [
            {
                "runs": [
                    {"text": "HLPFuzz", "bold": True, "color": GREEN_DARK, "size": 13},
                    {"text": "  · USENIX Sec'25 ", "italic": True, "color": GRAY_600, "size": 10},
                ]
            }
        ],
    )
    add_text(
        slide,
        right_x + Inches(0.2), by + Inches(0.46),
        right_w - Inches(0.3), card_h - Inches(0.5),
        [
            {"text": "种子生成：LLM 驱动的渐进式约束求解，选定未覆盖的基本块，LLM 生成覆盖该基本块的种子。", "space_after": 4},
            {"text": "预言机：常规的编译器崩溃，并不关心安全契约。"},
        ],
        font_size=10, color=GRAY_800, line_spacing=1.3,
    )

    # 底部 ribbon: oracle 空白
    rb_y = SLIDE_H - FOOTER_H - ribbon_h - Inches(0.04)
    add_rect(slide, CONTENT_PAD_X, rb_y, CONTENT_W, ribbon_h, fill=YELLOW_LIGHT, line=YELLOW_ACCENT, line_width=0.75)
    add_rect(slide, CONTENT_PAD_X, rb_y, Inches(0.08), ribbon_h, fill=YELLOW_ACCENT)
    add_text(
        slide,
        CONTENT_PAD_X + Inches(0.2), rb_y + Inches(0.05),
        CONTENT_W - Inches(0.3), Inches(0.28),
        ["Oracle 空白"],
        font_size=11, bold=True, color=GREEN_DARK,
    )
    add_text(
        slide,
        CONTENT_PAD_X + Inches(0.2), rb_y + Inches(0.32),
        CONTENT_W - Inches(0.3), ribbon_h - Inches(0.34),
        [
            {
                "runs": [
                    {"text": "现有的编译器崩溃、差分预言机面对变异", "size": 10},
                    {"text": "不报错、正常运行", "size": 10, "bold": True, "color": RED_DARK},
                    {"text": "的安全机制静默失效束手无策。", "size": 10},
                ]
            }
        ],
        font_size=10, color=GRAY_800, line_spacing=1.3,
    )
    return slide


# ═════════════════════════════════════════════════════════════
# 章节 03 — 研究目标与内容 (slides 5-6)
# ═════════════════════════════════════════════════════════════


def slide_section_03_divider(pres, pc):
    slide = _new_slide(pres)
    _render_section_divider(
        slide,
        number="03",
        title_zh="研究目标与内容",
        title_en="Goals & Content",
        footnote="核心目标：精准捕获“静默失效”；研究内容：构建从不变量到预言机的自动化检测链路。",
    )
    return slide


def slide_05_goals(pres, pc):
    """研究目标 (a/b/c) + 顶部 CVE-2023-4039 锚点 narrative (融合故事线)."""
    slide = _new_slide(pres)
    add_page_frame(slide, "研究目标与内容 — 研究目标", page_no=pc.next())

    y = CONTENT_TOP
    add_section_title(
        slide, CONTENT_PAD_X, y, CONTENT_W,
        "研究目标：填补防御机制的检测空白，捕获真实静默失效",
    )

    # ── 顶部锚点 / framing 卡 (融合故事线: 双重无声 + CVE-2023-4039 现实参考)
    intro_y = y + Inches(0.45)
    intro_h = Inches(0.85)
    add_rect(slide, CONTENT_PAD_X, intro_y, CONTENT_W, intro_h, fill=YELLOW_LIGHT, line=YELLOW_ACCENT, line_width=0.75)
    add_rect(slide, CONTENT_PAD_X, intro_y, Inches(0.08), intro_h, fill=YELLOW_ACCENT)
    add_text(
        slide,
        CONTENT_PAD_X + Inches(0.2), intro_y + Inches(0.08),
        CONTENT_W - Inches(0.3), Inches(0.3),
        [
            {
                "runs": [
                    {"text": "目标参照: ", "bold": True, "color": GREEN_DARK, "size": 11.5},
                    {"text": "CVE-2023-4039", "bold": True, "color": RED_DARK, "size": 11.5},
                    {"text": "  ·  AArch64 上 ", "color": GRAY_700, "size": 11.5, "italic": True},
                    {"text": "-fstack-protector", "color": GRAY_700, "size": 11.5, "italic": True},
                    {"text": " 形同虚设, 编译时无错、运行时正常, 但 canary 已被绕开", "color": GRAY_700, "size": 11.5, "italic": True},
                ]
            }
        ],
    )
    add_text(
        slide,
        CONTENT_PAD_X + Inches(0.2), intro_y + Inches(0.42),
        CONTENT_W - Inches(0.3), intro_h - Inches(0.46),
        [
            "我们以验证“安全契约”为目标，代替寻找程序崩溃。通过在“机制 × 架构”矩阵上系统排查，直接把这类“双重无声”的失效揪出来。"
        ],
        font_size=10.5, color=GRAY_800, line_spacing=1.3,
    )

    # ── 三大目标卡 (a/b/c)
    body_top = intro_y + intro_h + Inches(0.18)
    body_h = SLIDE_H - FOOTER_H - body_top - Inches(0.1)
    gap = Inches(0.18)
    card_w = (CONTENT_W - gap * 2) / 3

    goals = [
        {
            "tag": "G (a)",
            "title": "安全不变量集合",
            "bullets": [
                "覆盖主流防御机制 (canary · FORTIFY · IBT 等)",
                "每条不变量给出可机器判定的证据形式",
                "GCC 与 LLVM 共用同一描述结构",
            ],
        },
        {
            "tag": "G (b)",
            "title": "自动化检测原型",
            "bullets": [
                "同时面向 GCC 与 LLVM",
                "跨 ISA: x86-64 · AArch64 · RISC-V · LoongArch",
                "覆盖率驱动，专注防御相关代码",
                "结合 LLM 生成多样性种子",
            ],
        },
        {
            "tag": "G (c)",
            "title": "真实缺陷验证",
            "bullets": [
                "发现 CVE-2023-4039 类的静默失效",
                "推动上游确认与回归测试合入",
                "验证方法的跨编译器迁移性",
            ],
        },
    ]
    cx = CONTENT_PAD_X
    for g in goals:
        add_rect(slide, cx, body_top, card_w, body_h, fill=GREEN_LIGHT, line=GREEN_DARK, line_width=0.75)
        add_rect(slide, cx, body_top, card_w, Inches(0.05), fill=GREEN_DARK)
        add_text(
            slide,
            cx + Inches(0.18), body_top + Inches(0.13),
            card_w - Inches(0.36), Inches(0.26),
            [g["tag"]],
            font_size=11, bold=True, italic=True, color=YELLOW_ACCENT,
        )
        add_text(
            slide,
            cx + Inches(0.18), body_top + Inches(0.4),
            card_w - Inches(0.36), Inches(0.38),
            [g["title"]],
            font_size=15, bold=True, color=GREEN_DARK,
        )
        bullets = [{"text": "• " + b, "space_after": 5} for b in g["bullets"]]
        add_text(
            slide,
            cx + Inches(0.2), body_top + Inches(0.85),
            card_w - Inches(0.4), body_h - Inches(0.95),
            bullets,
            font_size=10.5, color=GRAY_800, line_spacing=1.3,
        )
        cx += card_w + gap
    return slide


def slide_06_content(pres, pc):
    """研究内容: 复用 pipeline 图 (上半) + 三项内容详细卡 (下半)."""
    slide = _new_slide(pres)
    add_page_frame(slide, "研究目标与内容 — 研究内容", page_no=pc.next())

    y = CONTENT_TOP
    add_section_title(
        slide, CONTENT_PAD_X, y, CONTENT_W,
        "研究内容：从安全契约到自动化模糊测试引擎",
    )

    body_top = y + Inches(0.45)

    # ── 上半: pipeline 图 (复用报告图)
    img_w = Inches(8.4)
    img_h = Inches(8.4 / 4.357)
    img_x = CONTENT_PAD_X + (CONTENT_W - img_w) / 2
    img_y = body_top
    slide.shapes.add_picture(str(IMG_PIPELINE), img_x, img_y, width=img_w, height=img_h)
    add_text(
        slide,
        img_x, img_y + img_h + Inches(0.02),
        img_w, Inches(0.2),
        ["图 4:  研究目标与内容"],
        font_size=9, italic=True, color=GRAY_600,
        align=PP_ALIGN.CENTER,
    )

    # ── 下半: 三项内容详细卡
    cards_top = img_y + img_h + Inches(0.32)
    cards_h = SLIDE_H - FOOTER_H - cards_top - Inches(0.1)
    gap = Inches(0.18)
    card_w = (CONTENT_W - gap * 2) / 3

    contents = [
        {
            "n": "①",
            "title": "提炼安全不变量（安全契约）",
            "color_fill": GREEN_LIGHT,
            "color_border": GREEN_DARK,
            "color_title": GREEN_DARK,
            "bullets": [
                "对主流防御机制分别进行调研，产出一套跨 ISA 的安全契约。",
            ],
        },
        {
            "n": "②",
            "title": "将不变量转化为可执行预言机",
            "color_fill": BLUE_LIGHT,
            "color_border": BLUE_DARK,
            "color_title": BLUE_DARK,
            "bullets": [
                "根据不变量的品相，实现为程序化的，静态检测或动态运行的检查器。",
                "任意检查器断言失败，即判定为漏洞。",
            ],
        },
        {
            "n": "③",
            "title": "LLM 与覆盖率驱动的主循环",
            "color_fill": YELLOW_LIGHT,
            "color_border": YELLOW_ACCENT,
            "color_title": GREEN_DARK,
            "bullets": [
                "通过防御机制相关代码的覆盖率来保证种子多样性",
                "将种子的多样性 与 预言机的健壮性 做乘法",
            ],
        },
    ]
    cx = CONTENT_PAD_X
    for c in contents:
        add_rect(slide, cx, cards_top, card_w, cards_h, fill=c["color_fill"], line=c["color_border"], line_width=0.75)
        add_rect(slide, cx, cards_top, card_w, Inches(0.05), fill=c["color_border"])
        add_text(
            slide,
            cx + Inches(0.16), cards_top + Inches(0.12),
            Inches(0.4), Inches(0.4),
            [c["n"]],
            font_size=20, bold=True, color=c["color_title"], anchor=MSO_ANCHOR.MIDDLE,
        )
        add_text(
            slide,
            cx + Inches(0.58), cards_top + Inches(0.16),
            card_w - Inches(0.7), Inches(0.36),
            [c["title"]],
            font_size=12, bold=True, color=c["color_title"], line_spacing=1.1,
        )
        bullets = [{"text": "• " + b, "space_after": 4} for b in c["bullets"]]
        add_text(
            slide,
            cx + Inches(0.18), cards_top + Inches(0.62),
            card_w - Inches(0.36), cards_h - Inches(0.7),
            bullets,
            font_size=9.5, color=GRAY_800, line_spacing=1.3,
        )
        cx += card_w + gap
    return slide


# ═════════════════════════════════════════════════════════════
# 章节 04 — 研究方案 (slides 8-10)
# ═════════════════════════════════════════════════════════════


def slide_section_04_divider(pres, pc):
    slide = _new_slide(pres)
    _render_section_divider(
        slide,
        number="04",
        title_zh="研究方案",
        title_en="Technical Route",
        footnote="从总体架构到核心算法：构建基于安全契约的自动化检测系统。",
    )
    return slide


def slide_08_method_arch(pres, pc):
    slide = _new_slide(pres)
    add_page_frame(slide, "研究方案 — 总体架构", page_no=pc.next())

    y = CONTENT_TOP
    add_section_title(
        slide, CONTENT_PAD_X, y, CONTENT_W,
        "总体架构：模块化设计与协同逻辑",
    )

    body_top = y + Inches(0.48)
    # 图比例 2.064: 宽景; 缩小以留出底部说明 ribbon 空间
    img_w = Inches(5.6)
    img_h = Inches(5.6 / 2.064)
    img_x = CONTENT_PAD_X + (CONTENT_W - img_w) / 2
    img_y = body_top
    slide.shapes.add_picture(str(IMG_ARCH), img_x, img_y, width=img_w, height=img_h)
    add_text(
        slide,
        img_x, img_y + img_h + Inches(0.04),
        img_w, Inches(0.22),
        ["图 5: 总体架构 — 子系统分工"],
        font_size=9.5, italic=True, color=GRAY_600,
        align=PP_ALIGN.CENTER,
    )

    # 底部说明 ribbon
    desc_y = img_y + img_h + Inches(0.34)
    desc_h = SLIDE_H - FOOTER_H - desc_y - Inches(0.08)
    add_rect(slide, CONTENT_PAD_X, desc_y, CONTENT_W, desc_h, fill=GRAY_100, line=GRAY_300, line_width=0.5)
    add_text(
        slide,
        CONTENT_PAD_X + Inches(0.18), desc_y + Inches(0.06),
        CONTENT_W - Inches(0.36), desc_h - Inches(0.12),
        [
            {"text": "• Fuzz 引擎: 主循环调度 + 提示词组装", "space_after": 2},
            {"text": "• 编译器子系统: 封装插桩 GCC / LLVM, 统一「行级覆盖 → 种子 ID」映射", "space_after": 2},
            {"text": "• 预言机子系统: 按机制聚合检查器, 输出安全判定; 配置模块按「机制 · ISA · flag」三元组下发任务"},
        ],
        font_size=10, color=GRAY_800, line_spacing=1.25,
    )
    return slide


def slide_09_method_oracle(pres, pc):
    slide = _new_slide(pres)
    add_page_frame(slide, "研究方案 — 预言机框架", page_no=pc.next())

    y = CONTENT_TOP
    add_section_title(
        slide, CONTENT_PAD_X, y, CONTENT_W,
        "预言机框架：静态与动态维度的交叉判定",
    )

    body_top = y + Inches(0.42)
    ribbon_h = Inches(0.5)
    body_h = CONTENT_H - Inches(0.5) - ribbon_h - Inches(0.16)

    # 左侧 oracle-flow 图 (比例 0.763 portrait)
    img_h = body_h
    img_w = int(body_h * 0.763)
    img_x = CONTENT_PAD_X + Inches(0.2)
    img_y = body_top
    slide.shapes.add_picture(str(IMG_ORACLE), img_x, img_y, width=img_w, height=img_h)
    add_text(
        slide,
        img_x, body_top + body_h + Inches(0.02),
        img_w, Inches(0.22),
        ["图 6: 预言机判定流程"],
        font_size=9.5, italic=True, color=GRAY_600,
        align=PP_ALIGN.CENTER,
    )

    # 右侧: 两类 checker 卡片 (上蓝 / 下绿)
    right_x = img_x + img_w + Inches(0.3)
    right_w = CONTENT_W + CONTENT_PAD_X - right_x
    card_gap = Inches(0.18)
    card_h = (body_h - card_gap) / 2

    # 静态检查器
    add_rect(slide, right_x, body_top, right_w, card_h, fill=BLUE_LIGHT, line=BLUE_DARK, line_width=0.75)
    add_rect(slide, right_x, body_top, Inches(0.08), card_h, fill=BLUE_DARK)
    add_text(
        slide,
        right_x + Inches(0.2), body_top + Inches(0.1),
        right_w - Inches(0.3), Inches(0.3),
        ["静态检查器"],
        font_size=13, bold=True, color=BLUE_DARK,
    )
    add_text(
        slide,
        right_x + Inches(0.2), body_top + Inches(0.5),
        right_w - Inches(0.3), card_h - Inches(0.55),
        [
            {"text": "• ELF 符号表检查 (如是否引用 __stack_chk_fail)", "space_after": 4},
            {"text": "• 指令流特征扫描 (如 endbr 指令位置)", "space_after": 4},
            {"text": "• .note 段属性检查"},
        ],
        font_size=10.5, color=GRAY_800, line_spacing=1.3,
    )

    # 动态检查器
    dy = body_top + card_h + card_gap
    add_rect(slide, right_x, dy, right_w, card_h, fill=GREEN_LIGHT, line=GREEN_DARK, line_width=0.75)
    add_rect(slide, right_x, dy, Inches(0.08), card_h, fill=GREEN_DARK)
    add_text(
        slide,
        right_x + Inches(0.2), dy + Inches(0.1),
        right_w - Inches(0.3), Inches(0.3),
        ["动态检查器"],
        font_size=13, bold=True, color=GREEN_DARK,
    )
    add_text(
        slide,
        right_x + Inches(0.2), dy + Inches(0.5),
        right_w - Inches(0.3), card_h - Inches(0.55),
        [
            {"text": "• 受控越界写入，观察能否被正常拦截", "space_after": 4},
            {"text": "• 植入汇编探针检测寄存器残留", "space_after": 4},
            {"text": "• 使用 QEMU-user 跨架构运行对照"},
        ],
        font_size=10.5, color=GRAY_800, line_spacing=1.3,
    )

    # 底部 ribbon: 判定语义 + 对照
    rb_y = SLIDE_H - FOOTER_H - ribbon_h - Inches(0.04)
    add_rect(slide, CONTENT_PAD_X, rb_y, CONTENT_W, ribbon_h, fill=YELLOW_LIGHT, line=YELLOW_ACCENT, line_width=0.75)
    add_rect(slide, CONTENT_PAD_X, rb_y, Inches(0.08), ribbon_h, fill=YELLOW_ACCENT)
    add_text(
        slide,
        CONTENT_PAD_X + Inches(0.2), rb_y + Inches(0.05),
        CONTENT_W - Inches(0.3), Inches(0.26),
        ["聚合判定与对照实验"],
        font_size=11, bold=True, color=GREEN_DARK,
    )
    add_text(
        slide,
        CONTENT_PAD_X + Inches(0.2), rb_y + Inches(0.28),
        CONTENT_W - Inches(0.3), ribbon_h - Inches(0.3),
        ["任何一处检查失败即视为安全契约被打破。系统会使用机制开/关两套编译参数进行对照，排除编译器缺省行为的干扰。"],
        font_size=10, color=GRAY_800, line_spacing=1.3,
    )
    return slide


def slide_10_method_loop(pres, pc):
    slide = _new_slide(pres)
    add_page_frame(slide, "研究方案 — LLM × 覆盖率主循环 + 跨架构调度", page_no=pc.next())

    y = CONTENT_TOP
    add_section_title(
        slide, CONTENT_PAD_X, y, CONTENT_W,
        "主循环：LLM 与覆盖率驱动的代码探索",
    )

    body_top = y + Inches(0.5)
    # 上半 3 个流程方框 (自绘, 不用 main-loop 图)
    flow_h = Inches(1.6)
    box_h = Inches(1.2)
    arrow_w = Inches(0.38)
    arrow_h = Inches(0.34)
    box_w = (CONTENT_W - arrow_w * 2) / 3
    bx = CONTENT_PAD_X
    by = body_top + (flow_h - box_h) / 2

    steps = [
        {
            "n": "1",
            "title": "选择目标",
            "sub": "控制流图 + 机制相关性",
            "detail": "优先挑选未覆盖且涉及防御逻辑的基本块",
        },
        {
            "n": "2",
            "title": "模型变异",
            "sub": "融合契约约束与代码上下文",
            "detail": "提供目标源码、覆盖标注与示例供模型参考",
        },
        {
            "n": "3",
            "title": "测试反馈",
            "sub": "验证覆盖率与安全性",
            "detail": "发现新覆盖或触发缺陷即保存；否则定位发散点重新变异",
        },
    ]
    for i, s in enumerate(steps):
        add_rect(slide, bx, by, box_w, box_h, fill=GREEN_LIGHT, line=GREEN_DARK, line_width=0.75)
        add_rect(slide, bx, by, box_w, Inches(0.05), fill=GREEN_DARK)
        # 圆形编号
        circle_d = Inches(0.46)
        add_rect(
            slide, bx + Inches(0.15), by + Inches(0.14),
            circle_d, circle_d,
            fill=GREEN_DARK, shape=MSO_SHAPE.OVAL,
        )
        add_text(
            slide,
            bx + Inches(0.15), by + Inches(0.14),
            circle_d, circle_d,
            [s["n"]],
            font_size=16, bold=True, color=WHITE,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )
        add_text(
            slide,
            bx + Inches(0.72), by + Inches(0.15),
            box_w - Inches(0.85), Inches(0.3),
            [s["title"]],
            font_size=13, bold=True, color=GREEN_DARK,
        )
        add_text(
            slide,
            bx + Inches(0.72), by + Inches(0.45),
            box_w - Inches(0.85), Inches(0.24),
            [s["sub"]],
            font_size=9.5, italic=True, color=GRAY_600,
        )
        add_text(
            slide,
            bx + Inches(0.2), by + Inches(0.76),
            box_w - Inches(0.4), box_h - Inches(0.82),
            [s["detail"]],
            font_size=10, color=GRAY_800, line_spacing=1.3,
        )
        if i < 2:
            add_arrow(
                slide,
                bx + box_w,
                by + (box_h - arrow_h) / 2,
                arrow_w, arrow_h,
                fill=GREEN_DARK,
            )
        bx += box_w + arrow_w

    # 下半: 跨架构调度 ribbon
    rb_y = body_top + flow_h + Inches(0.25)
    rb_h = SLIDE_H - FOOTER_H - rb_y - Inches(0.12)
    add_rect(slide, CONTENT_PAD_X, rb_y, CONTENT_W, rb_h, fill=YELLOW_LIGHT, line=YELLOW_ACCENT, line_width=0.75)
    add_rect(slide, CONTENT_PAD_X, rb_y, Inches(0.08), rb_h, fill=YELLOW_ACCENT)
    add_text(
        slide,
        CONTENT_PAD_X + Inches(0.2), rb_y + Inches(0.08),
        CONTENT_W - Inches(0.3), Inches(0.3),
        ["跨架构与编译选项调度"],
        font_size=12, bold=True, color=GREEN_DARK,
    )
    add_text(
        slide,
        CONTENT_PAD_X + Inches(0.2), rb_y + Inches(0.4),
        CONTENT_W - Inches(0.3), rb_h - Inches(0.45),
        [
            {"text": "• 按“机制 · 架构 · 参数”下发任务，同一颗种子跨多架构在 QEMU 中执行。", "space_after": 3},
            {"text": "• 采用机制开/关对照组，避免把编译器的正常行为误判为安全失效。"},
        ],
        font_size=10, color=GRAY_800, line_spacing=1.3,
    )
    return slide


# ═════════════════════════════════════════════════════════════
# 章节 05 — 现有成果 (slides 11-13)
# ═════════════════════════════════════════════════════════════


def slide_section_05_divider(pres, pc):
    slide = _new_slide(pres)
    _render_section_divider(
        slide,
        number="05",
        title_zh="现有成果",
        title_en="Preliminary Results",
        footnote="初步成果：已在 GCC 多个架构与版本中捕获 3 例静默失效缺陷。",
    )
    return slide


def slide_11_results_overview(pres, pc):
    slide = _new_slide(pres)
    add_page_frame(slide, "现有成果 — GCC 上的三例静默失效", page_no=pc.next())

    y = CONTENT_TOP
    add_section_title(
        slide, CONTENT_PAD_X, y, CONTENT_W,
        "初步成果：在 GCC 多个架构中捕获真实静默失效",
    )

    body_top = y + Inches(0.5)
    foot_h = Inches(0.42)
    body_h = CONTENT_H - Inches(0.55) - foot_h - Inches(0.1)

    gap = Inches(0.18)
    card_w = (CONTENT_W - gap * 2) / 3

    bugs = [
        {
            "tag": "Bug 1",
            "title": "_FORTIFY_SOURCE size 静默放宽",
            "where": "tree-ssa-strlen pass",
            "pr": "PR-125079",
            "status": "✓ 已合入 master",
            "status_color": GREEN_DARK,
            "desc": "将 __strcat_chk 优化为 __stpcpy_chk 时，未相应减小 size 参数。这导致运行时的安全上限被暗中放宽。",
        },
        {
            "tag": "Bug 2",
            "title": "IBT endbr 字节模式漏检",
            "where": "x86 后端 endbr 谓词",
            "pr": "PR-125084",
            "status": "⏳ 等待维护者回复",
            "status_color": "orange",
            "desc": "扫描循环遇到立即数高位非零时提前退出，嵌在高位的 endbr 字节被漏检。代码段中被混入了可被劫持的伪跳转入口。",
        },
        {
            "tag": "Bug 3",
            "title": "多 ISA canary 寄存器残留",
            "where": "通用 SSP fallback 后端",
            "pr": "PR-125045 (meta) · PR-125049",
            "status": "✓ LoongArch 子例已合入",
            "status_color": GREEN_DARK,
            "desc": "部分架构的后备实现在函数返回前未清零金丝雀寄存器，导致进程级的全局哨兵值残留在通用寄存器中。",
        },
    ]
    cx = CONTENT_PAD_X
    for b in bugs:
        add_rect(slide, cx, body_top, card_w, body_h, fill=RED_LIGHT, line=RED_DARK, line_width=0.75)
        add_rect(slide, cx, body_top, card_w, Inches(0.05), fill=RED_DARK)
        # 标签 (Bug N)
        add_text(
            slide,
            cx + Inches(0.18), body_top + Inches(0.14),
            card_w - Inches(0.36), Inches(0.26),
            [b["tag"]],
            font_size=11, bold=True, italic=True, color=YELLOW_ACCENT,
        )
        # 标题
        add_text(
            slide,
            cx + Inches(0.18), body_top + Inches(0.42),
            card_w - Inches(0.36), Inches(0.6),
            [b["title"]],
            font_size=13, bold=True, color=GRAY_800, line_spacing=1.1,
        )
        # where
        add_text(
            slide,
            cx + Inches(0.18), body_top + Inches(1.0),
            card_w - Inches(0.36), Inches(0.24),
            [
                {
                    "runs": [
                        {"text": "位置: ", "italic": True, "color": GRAY_600, "size": 9.5},
                        {"text": b["where"], "color": GRAY_800, "size": 9.5},
                    ]
                }
            ],
        )
        # PR
        add_text(
            slide,
            cx + Inches(0.18), body_top + Inches(1.24),
            card_w - Inches(0.36), Inches(0.24),
            [
                {
                    "runs": [
                        {"text": "上游: ", "italic": True, "color": GRAY_600, "size": 9.5},
                        {"text": b["pr"], "color": GRAY_800, "size": 9.5, "bold": True},
                    ]
                }
            ],
        )
        # 描述
        add_text(
            slide,
            cx + Inches(0.18), body_top + Inches(1.55),
            card_w - Inches(0.36), body_h - Inches(2.05),
            [b["desc"]],
            font_size=10, color=GRAY_800, line_spacing=1.3,
        )
        # 状态条 (底部一行)
        status_y = body_top + body_h - Inches(0.4)
        from theme import ORANGE_DARK  # local import to keep top-level imports tidy
        status_color = b["status_color"] if b["status_color"] != "orange" else ORANGE_DARK
        add_rect(slide, cx, status_y, card_w, Inches(0.4), fill=WHITE, line=status_color, line_width=0.5)
        add_text(
            slide,
            cx + Inches(0.12), status_y, card_w - Inches(0.24), Inches(0.4),
            [b["status"]],
            font_size=10.5, bold=True, color=status_color,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        cx += card_w + gap

    # 底部小字: 实验范围
    foot_y = SLIDE_H - FOOTER_H - foot_h - Inches(0.02)
    add_text(
        slide,
        CONTENT_PAD_X, foot_y, CONTENT_W, foot_h,
        [
            "当前实验范围涵盖 GCC 的多个架构与版本；"
            "针对 LLVM 及其他防御机制的测试将在下一阶段展开。"
        ],
        font_size=9, italic=True, color=GRAY_600, line_spacing=1.25,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    return slide


def slide_12_canary_principle(pres, pc):
    slide = _new_slide(pres)
    add_page_frame(slide, "现有成果 — Canary 寄存器残留: 原理与历史教训", page_no=pc.next())

    y = CONTENT_TOP
    add_section_title(
        slide, CONTENT_PAD_X, y, CONTENT_W,
        "缺陷剖析：Canary 寄存器残留的破坏性",
    )

    body_top = y + Inches(0.5)
    foot_h = Inches(0.5)
    body_h = CONTENT_H - Inches(0.55) - foot_h - Inches(0.1)

    gap = Inches(0.22)
    col_w = (CONTENT_W - gap) / 2

    # 左: Canary 原理
    add_rect(slide, CONTENT_PAD_X, body_top, col_w, body_h, fill=GREEN_LIGHT, line=GREEN_DARK, line_width=0.75)
    add_rect(slide, CONTENT_PAD_X, body_top, col_w, Inches(0.05), fill=GREEN_DARK)
    add_text(
        slide,
        CONTENT_PAD_X + Inches(0.18), body_top + Inches(0.14),
        col_w - Inches(0.36), Inches(0.34),
        ["安全规范：不能让底牌漏出"],
        font_size=13, bold=True, color=GREEN_DARK,
    )
    add_text(
        slide,
        CONTENT_PAD_X + Inches(0.2), body_top + Inches(0.55),
        col_w - Inches(0.4), body_h - Inches(0.65),
        [
            {
                "runs": [
                    {"text": "① Canary 值来自全局 ", "size": 10.5, "color": GRAY_800},
                    {"text": "__stack_chk_guard", "size": 10.5, "color": GREEN_DARK, "bold": True},
                    {"text": " 并在进程内复用。一旦这个值被窃取，后续防线就会彻底瓦解。", "size": 10.5, "color": GRAY_800},
                ],
                "space_after": 8,
            },
            {
                "runs": [
                    {"text": "② Canary 需要频繁加载到通用寄存器中进行校验或压栈。", "size": 10.5, "color": GRAY_800},
                ],
                "space_after": 8,
            },
            {
                "runs": [
                    {"text": "③ 安全规范强制要求：", "size": 10.5, "color": GRAY_800},
                    {"text": "函数返回前必须把用过的金丝雀寄存器清零", "size": 10.5, "color": GREEN_DARK, "bold": True},
                    {"text": "。否则该值就会随寄存器泄漏给攻击者。", "size": 10.5, "color": GRAY_800},
                ],
            },
        ],
        line_spacing=1.35,
    )

    # 右: 历史教训
    rx = CONTENT_PAD_X + col_w + gap
    add_rect(slide, rx, body_top, col_w, body_h, fill=YELLOW_LIGHT, line=YELLOW_ACCENT, line_width=0.75)
    add_rect(slide, rx, body_top, col_w, Inches(0.05), fill=YELLOW_ACCENT)
    add_text(
        slide,
        rx + Inches(0.18), body_top + Inches(0.14),
        col_w - Inches(0.36), Inches(0.34),
        ["历史重演：同源缺陷跨架构潜伏"],
        font_size=13, bold=True, color=GREEN_DARK,
    )
    add_text(
        slide,
        rx + Inches(0.2), body_top + Inches(0.55),
        col_w - Inches(0.4), body_h - Inches(0.65),
        [
            {
                "runs": [
                    {"text": "2020 年，GCC 曾修复过 ARM 架构上的金丝雀寄存器残留问题 (", "size": 10.5, "color": GRAY_800},
                    {"text": "PR-96191", "size": 10.5, "color": GREEN_DARK, "bold": True},
                    {"text": ")。", "size": 10.5, "color": GRAY_800},
                ],
                "space_after": 10,
            },
            {
                "runs": [
                    {"text": "但是，当时并未同步修复通用的后备实现。MIPS、LoongArch 等架构一直使用的是这套有缺陷的后备代码。", "size": 10.5, "color": GRAY_800},
                ],
                "space_after": 10,
            },
            {
                "runs": [
                    {"text": "由于程序不会崩溃，这个遗漏在 5 年内无人察觉。系统看似开启了保护，", "size": 10.5, "color": GRAY_800},
                    {"text": "金丝雀值却随时可能被读取", "size": 10.5, "color": RED_DARK, "bold": True},
                    {"text": "。", "size": 10.5, "color": GRAY_800},
                ],
            },
        ],
        line_spacing=1.35,
    )

    # 底部 ribbon
    rb_y = SLIDE_H - FOOTER_H - foot_h - Inches(0.04)
    add_rect(slide, CONTENT_PAD_X, rb_y, CONTENT_W, foot_h, fill=GRAY_100, line=GRAY_300, line_width=0.5)
    add_text(
        slide,
        CONTENT_PAD_X + Inches(0.18), rb_y + Inches(0.06),
        CONTENT_W - Inches(0.36), foot_h - Inches(0.12),
        [
            {
                "runs": [
                    {"text": "我们的解法：", "bold": True, "color": GREEN_DARK, "size": 11},
                    {"text": "将“函数返回前必须清零金丝雀寄存器”提炼为一条安全不变量，并用动态探针加以验证。", "color": GRAY_800, "size": 11},
                ]
            }
        ],
    )
    return slide


def slide_13_canary_checker(pres, pc):
    slide = _new_slide(pres)
    add_page_frame(slide, "现有成果 — 用内联汇编 checker 触发多 ISA 泄漏", page_no=pc.next())

    y = CONTENT_TOP
    add_section_title(
        slide, CONTENT_PAD_X, y, CONTENT_W,
        "实证发现：内联汇编探针触发架构级泄露",
    )

    body_top = y + Inches(0.48)
    foot_h = Inches(0.72)
    body_h = CONTENT_H - Inches(0.53) - foot_h - Inches(0.1)

    gap = Inches(0.22)
    left_w = Inches(4.5)
    right_w = CONTENT_W - left_w - gap

    # 左: checker 设计 + 短例
    add_rect(slide, CONTENT_PAD_X, body_top, left_w, body_h, fill=GREEN_LIGHT, line=GREEN_DARK, line_width=0.75)
    add_rect(slide, CONTENT_PAD_X, body_top, left_w, Inches(0.05), fill=GREEN_DARK)
    add_text(
        slide,
        CONTENT_PAD_X + Inches(0.18), body_top + Inches(0.14),
        left_w - Inches(0.36), Inches(0.34),
        ["动态探针设计"],
        font_size=13, bold=True, color=GREEN_DARK,
    )
    add_text(
        slide,
        CONTENT_PAD_X + Inches(0.2), body_top + Inches(0.55),
        left_w - Inches(0.4), Inches(1.1),
        [
            {
                "runs": [
                    {"text": "在函数返回前插入汇编指令，读取各个寄存器的值并与 ", "size": 10.5, "color": GRAY_800},
                    {"text": "__stack_chk_guard", "size": 10.5, "color": GREEN_DARK, "bold": True},
                    {"text": " 对比；只要发现一致，即可确认 ", "size": 10.5, "color": GRAY_800},
                    {"text": "金丝雀发生了泄漏", "size": 10.5, "color": RED_DARK, "bold": True},
                    {"text": "。", "size": 10.5, "color": GRAY_800},
                ],
            }
        ],
        line_spacing=1.3,
    )
    # 代码示例 (LoongArch)
    code_y = body_top + Inches(1.78)
    code_h = body_h - Inches(1.85)
    add_rect(slide, CONTENT_PAD_X + Inches(0.18), code_y, left_w - Inches(0.36), code_h, fill=GRAY_100, line=GRAY_300, line_width=0.5)
    add_text(
        slide,
        CONTENT_PAD_X + Inches(0.28), code_y + Inches(0.06),
        left_w - Inches(0.56), Inches(0.24),
        ["LoongArch 泄漏探针示例"],
        font_size=9, italic=True, color=GRAY_600,
    )
    add_text(
        slide,
        CONTENT_PAD_X + Inches(0.28), code_y + Inches(0.3),
        left_w - Inches(0.56), code_h - Inches(0.36),
        [
            {"text": "asm volatile (", "space_after": 0},
            {"text": "    \"la.local $t1, __stack_chk_guard\\n\"", "space_after": 0},
            {"text": "    \"ld.d   $t1, $t1, 0\\n\"", "space_after": 0},
            {"text": "    \"beq    $t1, $r12, leaked\\n\"", "space_after": 0},
            {"text": ");  // 返回前若 $r12 还存着 guard 值 → 泄漏", "space_after": 0},
        ],
        font_size=8.5, color=GRAY_800,
        font_name="Courier New", ea_name="Microsoft YaHei",
        line_spacing=1.25,
    )

    # 右: 截图
    rx = CONTENT_PAD_X + left_w + gap
    # 图比例 1.715, 限定宽度 right_w
    img_w = right_w
    img_h = int(right_w * (1.0 / 1.715))
    if img_h > body_h - Inches(0.35):
        img_h = body_h - Inches(0.35)
        img_w = int(img_h * 1.715)
    img_x = rx + (right_w - img_w) / 2
    img_y = body_top + (body_h - img_h - Inches(0.28)) / 2
    slide.shapes.add_picture(str(IMG_GCC_DEV), img_x, img_y, width=img_w, height=img_h)
    add_text(
        slide,
        rx, img_y + img_h + Inches(0.04),
        right_w, Inches(0.24),
        ["图 7: 上游对该安全隐患的确认邮件"],
        font_size=9.5, italic=True, color=GRAY_600,
        align=PP_ALIGN.CENTER,
    )

    # 底部 ribbon: 后续
    rb_y = SLIDE_H - FOOTER_H - foot_h - Inches(0.04)
    add_rect(slide, CONTENT_PAD_X, rb_y, CONTENT_W, foot_h, fill=GREEN_LIGHT, line=GREEN_DARK, line_width=0.75)
    add_rect(slide, CONTENT_PAD_X, rb_y, Inches(0.08), foot_h, fill=GREEN_DARK)
    add_text(
        slide,
        CONTENT_PAD_X + Inches(0.2), rb_y + Inches(0.06),
        CONTENT_W - Inches(0.3), Inches(0.28),
        ["上游反馈"],
        font_size=11, bold=True, color=GREEN_DARK,
    )
    add_text(
        slide,
        CONTENT_PAD_X + Inches(0.2), rb_y + Inches(0.32),
        CONTENT_W - Inches(0.3), foot_h - Inches(0.35),
        [
            {
                "runs": [
                    {"text": "LoongArch 维护者将其定性为严重的系统漏洞 (", "color": GRAY_800, "size": 10},
                    {"text": "severe security vulnerability", "bold": True, "color": RED_DARK, "size": 10},
                    {"text": ")，并决定将修复移植到旧版本。其他架构的同类隐患也正在依次上报。", "color": GRAY_800, "size": 10},
                ]
            }
        ],
        line_spacing=1.3,
    )
    return slide


# ═════════════════════════════════════════════════════════════
# 封面
# ═════════════════════════════════════════════════════════════


def slide_cover(pres, pc):
    """封面: 校徽 + 论文题目 + 作者信息. 不占页码."""
    slide = _new_slide(pres)

    # 顶部绿色装饰条 + 底部黄色细条
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.18), fill=GREEN_DARK)
    add_rect(slide, 0, SLIDE_H - Inches(0.08), SLIDE_W, Inches(0.08), fill=YELLOW_ACCENT)

    # 校徽 (居中略偏上)
    logo_w = Inches(0.95)
    logo_h = Inches(0.95)
    logo_x = (SLIDE_W - logo_w) / 2
    logo_y = Inches(0.48)
    slide.shapes.add_picture(str(LOGO_PATH), logo_x, logo_y, width=logo_w, height=logo_h)

    # 学校 / 院系
    add_text(
        slide,
        Inches(0.5), logo_y + logo_h + Inches(0.05),
        SLIDE_W - Inches(1.0), Inches(0.3),
        ["东南大学  ·  网络空间安全学院"],
        font_size=14, color=GREEN_DARK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )

    # 论文标题 (主)
    add_text(
        slide,
        Inches(0.5), Inches(2.05),
        SLIDE_W - Inches(1.0), Inches(0.7),
        ["基于大语言模型的编译器防御机制脆弱性分析方法研究"],
        font_size=28, bold=True, color=GREEN_DARK,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )
    # 黄色短分隔线
    bar_w = Inches(0.6)
    add_rect(
        slide,
        (SLIDE_W - bar_w) / 2, Inches(2.78),
        bar_w, Inches(0.05),
        fill=YELLOW_ACCENT,
    )
    # 副标题 / 项目代号
    add_text(
        slide,
        Inches(0.5), Inches(2.9),
        SLIDE_W - Inches(1.0), Inches(0.36),
        ["DeFuzz · LLM 驱动的编译器防御机制模糊测试"],
        font_size=14, italic=True, color=GRAY_700,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )

    # 作者信息卡 (居中, 浅绿底)
    info_w = Inches(5.6)
    info_h = Inches(1.25)
    info_x = (SLIDE_W - info_w) / 2
    info_y = Inches(3.55)
    add_rect(slide, info_x, info_y, info_w, info_h, fill=GREEN_LIGHT, line=GREEN_DARK, line_width=0.75)
    add_rect(slide, info_x, info_y, Inches(0.08), info_h, fill=GREEN_DARK)

    # 左列: 研究生 / 学号
    col1_x = info_x + Inches(0.32)
    col1_w = Inches(2.5)
    add_text(
        slide,
        col1_x, info_y + Inches(0.18),
        col1_w, Inches(0.4),
        [
            {
                "runs": [
                    {"text": "研究生：", "color": GRAY_600, "size": 12},
                    {"text": "张景耀", "bold": True, "color": GREEN_DARK, "size": 14},
                ]
            }
        ],
    )
    add_text(
        slide,
        col1_x, info_y + Inches(0.6),
        col1_w, Inches(0.4),
        [
            {
                "runs": [
                    {"text": "学　号：", "color": GRAY_600, "size": 12},
                    {"text": "245569", "bold": True, "color": GRAY_800, "size": 13},
                ]
            }
        ],
    )

    # 右列: 导师 / 学位类别 / 专业领域
    col2_x = info_x + info_w / 2 + Inches(0.05)
    col2_w = info_w / 2 - Inches(0.25)
    add_text(
        slide,
        col2_x, info_y + Inches(0.18),
        col2_w, Inches(0.4),
        [
            {
                "runs": [
                    {"text": "指导教师：", "color": GRAY_600, "size": 12},
                    {"text": "徐坚皓", "bold": True, "color": GREEN_DARK, "size": 14},
                ]
            }
        ],
    )
    add_text(
        slide,
        col2_x, info_y + Inches(0.6),
        col2_w, Inches(0.4),
        [
            {
                "runs": [
                    {"text": "学位类别：", "color": GRAY_600, "size": 12},
                    {"text": "专业硕士", "bold": True, "color": GRAY_800, "size": 13},
                    {"text": "  ·  电子信息", "color": GRAY_700, "size": 12},
                ]
            }
        ],
    )

    # 底部日期
    add_text(
        slide,
        Inches(0.5), SLIDE_H - Inches(0.55),
        SLIDE_W - Inches(1.0), Inches(0.3),
        ["开题报告  ·  2026.05"],
        font_size=11, italic=True, color=GRAY_700,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )
    return slide


# ═════════════════════════════════════════════════════════════
# 章节 06 — 进度安排
# ═════════════════════════════════════════════════════════════


def slide_section_06_divider(pres, pc):
    slide = _new_slide(pres)
    _render_section_divider(
        slide,
        number="06",
        title_zh="进度安排",
        title_en="Schedule & Plan",
        footnote="从扩展机制到上游闭环：四阶段推进，保障按期完成学位论文。",
    )
    return slide


def slide_14_plan(pres, pc):
    """学位论文工作进度与安排: 时间轴卡片网格."""
    slide = _new_slide(pres)
    add_page_frame(slide, "进度安排 — 学位论文工作计划", page_no=pc.next())

    y = CONTENT_TOP
    add_section_title(
        slide, CONTENT_PAD_X, y, CONTENT_W,
        "工作进度：四阶段推进，确保按期答辩",
    )

    body_top = y + Inches(0.45)
    body_h = SLIDE_H - FOOTER_H - body_top - Inches(0.12)

    items = [
        ("2026.05 — 2026.07", "扩展防御机制覆盖",
         "在 canary、_FORTIFY_SOURCE、IBT 基础上，扩展到 BTI / PAC、Shadow Stack、CFI 等机制；丰富预言机的静态与动态检查器。"),
        ("2026.07 — 2026.09", "迁移到 LLVM 工具链",
         "把编译器子系统、覆盖率适配层、flag profile 扩展到 LLVM；对齐跨架构 / 跨编译器实验框架。"),
        ("2026.09 — 2026.12", "全机制端到端实验",
         "在 GCC 与 LLVM 多版本、多 ISA 上完成全机制端到端实验，持续推进上游缺陷上报闭环。"),
        ("2026.12 — 2027.02", "基线对照评测",
         "与 Csmith、YARPGen、GrayC、HLPFuzz 等基线方法对照，量化覆盖率推进与缺陷发现上的相对优势。"),
        ("2027.02 — 2027.04", "论文初稿与投稿",
         "整理实验数据与图表，完成学位论文初稿，准备会议 / 期刊投稿。"),
        ("2027.04 — 2027.06", "修订与答辩准备",
         "学位论文修订、外审反馈处理、答辩材料准备；2027.06 学位论文答辩。"),
    ]

    cols = 3
    rows = 2
    gap_x = Inches(0.18)
    gap_y = Inches(0.18)
    card_w = (CONTENT_W - gap_x * (cols - 1)) / cols
    card_h = (body_h - gap_y * (rows - 1)) / rows

    for i, (period, title, desc) in enumerate(items):
        r, c = divmod(i, cols)
        cx = CONTENT_PAD_X + c * (card_w + gap_x)
        cy = body_top + r * (card_h + gap_y)

        add_rect(slide, cx, cy, card_w, card_h, fill=GREEN_LIGHT, line=GREEN_DARK, line_width=0.75)
        add_rect(slide, cx, cy, card_w, Inches(0.05), fill=GREEN_DARK)

        # 阶段标签 (S1..S6)
        add_text(
            slide,
            cx + Inches(0.18), cy + Inches(0.1),
            card_w - Inches(0.36), Inches(0.26),
            [f"S{i + 1}"],
            font_size=10, bold=True, italic=True, color=YELLOW_ACCENT,
        )
        # 时间区间
        add_text(
            slide,
            cx + Inches(0.18), cy + Inches(0.3),
            card_w - Inches(0.36), Inches(0.28),
            [period],
            font_size=11, bold=True, color=GREEN_DARK,
        )
        # 标题
        add_text(
            slide,
            cx + Inches(0.18), cy + Inches(0.62),
            card_w - Inches(0.36), Inches(0.32),
            [title],
            font_size=12, bold=True, color=GRAY_800,
        )
        # 内容
        add_text(
            slide,
            cx + Inches(0.18), cy + Inches(0.98),
            card_w - Inches(0.36), card_h - Inches(1.05),
            [desc],
            font_size=9.5, color=GRAY_800, line_spacing=1.3,
        )
    return slide


# ═════════════════════════════════════════════════════════════
# 致谢页
# ═════════════════════════════════════════════════════════════


def slide_thanks(pres, pc):
    """结尾致谢页. 不占页码."""
    slide = _new_slide(pres)

    # 顶绿条 + 底黄条
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.18), fill=GREEN_DARK)
    add_rect(slide, 0, SLIDE_H - Inches(0.08), SLIDE_W, Inches(0.08), fill=YELLOW_ACCENT)

    # 大字 "感谢聆听"
    add_text(
        slide,
        Inches(0.5), Inches(1.7),
        SLIDE_W - Inches(1.0), Inches(1.3),
        ["感谢聆听"],
        font_size=72, bold=True, color=GREEN_DARK,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )
    # 黄色短线
    bar_w = Inches(1.0)
    add_rect(
        slide,
        (SLIDE_W - bar_w) / 2, Inches(3.05),
        bar_w, Inches(0.06),
        fill=YELLOW_ACCENT,
    )
    # 副标题
    add_text(
        slide,
        Inches(0.5), Inches(3.2),
        SLIDE_W - Inches(1.0), Inches(0.5),
        ["敬请各位老师批评指正"],
        font_size=20, color=GRAY_700,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )
    # 落款
    add_text(
        slide,
        Inches(0.5), SLIDE_H - Inches(0.7),
        SLIDE_W - Inches(1.0), Inches(0.32),
        ["张景耀  ·  指导教师：徐坚皓  ·  东南大学网络空间安全学院  ·  2026.05"],
        font_size=11, italic=True, color=GRAY_700,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )
    return slide


# ═════════════════════════════════════════════════════════════
# Deck slide 顺序 — 自动维护页码
# ═════════════════════════════════════════════════════════════
ALL_SLIDES = [
    # 封面
    slide_cover,
    # 章节 01
    slide_section_01_divider,
    slide_01_cve_growth,
    slide_02_canary_aarch64,
    slide_03_core_challenge,
    # 章节 02
    slide_section_02_divider,
    slide_04_status_methods,
    # 章节 03
    slide_section_03_divider,
    slide_06_content,
    # 章节 04
    slide_section_04_divider,
    slide_08_method_arch,
    slide_09_method_oracle,
    # 章节 05
    slide_section_05_divider,
    slide_11_results_overview,
    slide_12_canary_principle,
    slide_13_canary_checker,
    # 章节 06 — 进度安排
    slide_section_06_divider,
    slide_14_plan,
    # 致谢
    slide_thanks,
]
